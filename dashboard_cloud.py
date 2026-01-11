import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import json
import sqlite3
import subprocess
import os
from pathlib import Path
from datetime import datetime
from io import BytesIO

# PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Configuração da página
st.set_page_config(
    page_title="Vagas Colégio Elo",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

BASE_DIR = Path(__file__).parent

# ============================================================
# SIDEBAR - MENU LATERAL
# ============================================================
with st.sidebar:
    st.markdown("""
        <div style='text-align: center; padding: 1rem 0;'>
            <h2 style='color: #ffffff; margin: 0;'>🎓 SIGA Vagas</h2>
            <p style='color: #64748b; font-size: 0.8rem;'>Colégio Elo</p>
        </div>
    """, unsafe_allow_html=True)

    st.markdown("---")

    # Toggle Tema (Dark/Light)
    tema_escuro = st.toggle("🌙 Modo Escuro", value=True, key="tema_toggle")

    st.markdown("---")

    # Navegação rápida
    st.markdown("### 🧭 Navegação")

    st.markdown("""
        <style>
            .nav-link {
                display: block;
                padding: 0.5rem 1rem;
                margin: 0.25rem 0;
                background: rgba(59, 130, 246, 0.1);
                border-radius: 8px;
                color: #94a3b8;
                text-decoration: none;
                transition: all 0.2s;
            }
            .nav-link:hover {
                background: rgba(59, 130, 246, 0.2);
                color: #ffffff;
            }
        </style>
    """, unsafe_allow_html=True)

    nav_items = [
        ("📊", "Visão Geral", "Métricas principais"),
        ("🎯", "Gauge & Treemap", "Ocupação visual"),
        ("📈", "Analytics", "Ranking e projeções"),
        ("🏫", "Unidades", "Nível 1"),
        ("📚", "Segmentos", "Nível 2"),
        ("📖", "Séries", "Nível 3"),
        ("🎓", "Turmas", "Nível 4"),
    ]

    for icon, titulo, desc in nav_items:
        st.markdown(f"""
            <div style='background: rgba(59, 130, 246, 0.1); border-radius: 8px; padding: 0.5rem 0.75rem; margin: 0.3rem 0; cursor: pointer;'>
                <span style='font-size: 1rem;'>{icon}</span>
                <span style='color: #ffffff; margin-left: 0.5rem;'>{titulo}</span>
                <div style='color: #64748b; font-size: 0.7rem; margin-left: 1.75rem;'>{desc}</div>
            </div>
        """, unsafe_allow_html=True)

    st.markdown("---")

    # Configuração de Alertas
    st.markdown("### 🔔 Alertas")

    with st.expander("⚙️ Configurar Limites"):
        alerta_critico = st.slider("❄️ Crítico (abaixo de)", 0, 100, 50, 5, key="alerta_critico")
        alerta_atencao = st.slider("⚠️ Atenção (abaixo de)", 0, 100, 70, 5, key="alerta_atencao")
        alerta_lotado = st.slider("🔥 Quase Lotado (acima de)", 0, 100, 95, 5, key="alerta_lotado")

        st.markdown(f"""
            <div style='font-size: 0.75rem; color: #64748b; margin-top: 0.5rem;'>
                <p>❄️ Crítico: &lt;{alerta_critico}%</p>
                <p>⚠️ Atenção: {alerta_critico}-{alerta_atencao}%</p>
                <p>🔥 Lotado: &gt;{alerta_lotado}%</p>
            </div>
        """, unsafe_allow_html=True)

    st.markdown("---")

    # Relatórios
    st.markdown("### 📄 Relatórios")

    tipo_relatorio = st.selectbox(
        "Tipo de Relatório",
        ["Resumo Executivo", "Detalhado por Unidade", "Análise de Tendências", "Turmas Críticas"],
        key="tipo_relatorio"
    )

    formato_export = st.radio(
        "Formato",
        ["PDF", "PowerPoint", "Excel"],
        horizontal=True,
        key="formato_export"
    )

    if st.button("📥 Gerar Relatório", use_container_width=True, key="btn_gerar_relatorio"):
        st.session_state['gerar_relatorio'] = True
        st.session_state['tipo_relatorio_selecionado'] = tipo_relatorio
        st.session_state['formato_export_selecionado'] = formato_export

    st.markdown("---")

    # Agendamento de Relatórios
    st.markdown("### ⏰ Agendamento")

    # Carrega configuração de agendamento existente
    schedule_config_path = BASE_DIR / "output" / "schedule_config.json"
    schedule_config = {}
    if schedule_config_path.exists():
        try:
            with open(schedule_config_path) as f:
                schedule_config = json.load(f)
        except:
            schedule_config = {}

    with st.expander("📅 Configurar Envio Automático", expanded=False):
        # Ativar/Desativar
        agendamento_ativo = st.toggle(
            "Ativar envio automático",
            value=schedule_config.get('ativo', False),
            key="agendamento_ativo"
        )

        # Frequência
        frequencia = st.selectbox(
            "Frequência",
            ["Diário", "Semanal", "Mensal"],
            index=["Diário", "Semanal", "Mensal"].index(schedule_config.get('frequencia', 'Diário')),
            key="agendamento_frequencia"
        )

        # Dia da semana (se semanal)
        dias_semana = ["Segunda", "Terça", "Quarta", "Quinta", "Sexta", "Sábado", "Domingo"]
        if frequencia == "Semanal":
            dia_semana = st.selectbox(
                "Dia da semana",
                dias_semana,
                index=schedule_config.get('dia_semana', 0),
                key="agendamento_dia_semana"
            )
        else:
            dia_semana = "Segunda"

        # Dia do mês (se mensal)
        if frequencia == "Mensal":
            dia_mes = st.number_input(
                "Dia do mês",
                min_value=1,
                max_value=28,
                value=schedule_config.get('dia_mes', 1),
                key="agendamento_dia_mes"
            )
        else:
            dia_mes = 1

        # Horário
        col_hora, col_min = st.columns(2)
        with col_hora:
            hora = st.selectbox(
                "Hora",
                list(range(0, 24)),
                index=schedule_config.get('hora', 6),
                key="agendamento_hora"
            )
        with col_min:
            minuto = st.selectbox(
                "Minuto",
                [0, 15, 30, 45],
                index=[0, 15, 30, 45].index(schedule_config.get('minuto', 0)) if schedule_config.get('minuto', 0) in [0, 15, 30, 45] else 0,
                key="agendamento_minuto"
            )

        # Tipo de relatório
        tipo_agendado = st.selectbox(
            "Tipo de Relatório",
            ["Resumo Executivo", "Detalhado por Unidade", "Turmas Críticas"],
            index=["Resumo Executivo", "Detalhado por Unidade", "Turmas Críticas"].index(
                schedule_config.get('tipo_relatorio', 'Resumo Executivo')
            ) if schedule_config.get('tipo_relatorio', 'Resumo Executivo') in ["Resumo Executivo", "Detalhado por Unidade", "Turmas Críticas"] else 0,
            key="agendamento_tipo"
        )

        # Canais de envio
        st.markdown("**Canais de envio:**")
        enviar_email_check = st.checkbox(
            "📧 Email",
            value=schedule_config.get('enviar_email', True),
            key="agendamento_email"
        )
        enviar_whatsapp_check = st.checkbox(
            "📱 WhatsApp",
            value=schedule_config.get('enviar_whatsapp', False),
            key="agendamento_whatsapp"
        )

        # Botão salvar
        if st.button("💾 Salvar Agendamento", use_container_width=True, key="btn_salvar_agendamento"):
            nova_config = {
                'ativo': agendamento_ativo,
                'frequencia': frequencia,
                'dia_semana': dias_semana.index(dia_semana) if frequencia == "Semanal" else 0,
                'dia_mes': dia_mes if frequencia == "Mensal" else 1,
                'hora': hora,
                'minuto': minuto,
                'tipo_relatorio': tipo_agendado,
                'enviar_email': enviar_email_check,
                'enviar_whatsapp': enviar_whatsapp_check,
                'ultima_atualizacao': datetime.now().isoformat()
            }

            try:
                with open(schedule_config_path, 'w') as f:
                    json.dump(nova_config, f, indent=2, ensure_ascii=False)
                st.success("✅ Agendamento salvo!")

                # Mostra próximo envio
                if agendamento_ativo:
                    st.info(f"📅 Próximo envio: {frequencia} às {hora:02d}:{minuto:02d}")
            except Exception as e:
                st.error(f"❌ Erro ao salvar: {e}")

        # Status atual
        if schedule_config.get('ativo', False):
            freq = schedule_config.get('frequencia', 'Diário')
            h = schedule_config.get('hora', 6)
            m = schedule_config.get('minuto', 0)

            status_text = f"🟢 Ativo • {freq} às {h:02d}:{m:02d}"
            if freq == "Semanal":
                dias = ["Seg", "Ter", "Qua", "Qui", "Sex", "Sáb", "Dom"]
                status_text += f" ({dias[schedule_config.get('dia_semana', 0)]})"
            elif freq == "Mensal":
                status_text += f" (dia {schedule_config.get('dia_mes', 1)})"

            st.markdown(f"""
                <div style='background: rgba(34, 197, 94, 0.1); border: 1px solid rgba(34, 197, 94, 0.3);
                            border-radius: 8px; padding: 0.5rem; text-align: center; margin-top: 0.5rem;'>
                    <span style='color: #22c55e; font-size: 0.8rem;'>{status_text}</span>
                </div>
            """, unsafe_allow_html=True)
        else:
            st.markdown("""
                <div style='background: rgba(100, 116, 139, 0.1); border: 1px solid rgba(100, 116, 139, 0.3);
                            border-radius: 8px; padding: 0.5rem; text-align: center; margin-top: 0.5rem;'>
                    <span style='color: #64748b; font-size: 0.8rem;'>🔴 Agendamento desativado</span>
                </div>
            """, unsafe_allow_html=True)

    st.markdown("---")

    # Info
    st.markdown("""
        <div style='text-align: center; color: #64748b; font-size: 0.7rem; padding: 1rem 0;'>
            <p>Versão 2.2</p>
            <p>Atualizado automaticamente às 6h</p>
        </div>
    """, unsafe_allow_html=True)

# CSS Corporate SaaS - Tema Dinâmico
if tema_escuro:
    # Dark Mode - Navy Blue
    bg_gradient = "linear-gradient(180deg, #0a1628 0%, #0f2137 50%, #132743 100%)"
    card_bg = "linear-gradient(145deg, #0d1f35 0%, #142d4c 100%)"
    text_color = "#ffffff"
    text_muted = "#94a3b8"
    border_color = "rgba(59, 130, 246, 0.2)"
    grid_color = "rgba(59, 130, 246, 0.1)"
else:
    # Light Mode
    bg_gradient = "linear-gradient(180deg, #f8fafc 0%, #e2e8f0 50%, #cbd5e1 100%)"
    card_bg = "linear-gradient(145deg, #ffffff 0%, #f1f5f9 100%)"
    text_color = "#1e293b"
    text_muted = "#64748b"
    border_color = "rgba(59, 130, 246, 0.3)"
    grid_color = "rgba(59, 130, 246, 0.15)"

st.markdown(f"""
<style>
    /* Tema base */
    .stApp {{
        background: {bg_gradient};
    }}

    /* Main container */
    .main .block-container {{
        padding: 2rem 3rem;
        max-width: 1400px;
    }}

    /* Headers */
    h1, h2, h3 {{
        color: {text_color} !important;
        font-weight: 600 !important;
    }}

    h1 {{
        font-size: 2.5rem !important;
        color: {text_color} !important;
    }}

    /* Metric cards */
    [data-testid="stMetric"] {{
        background: {card_bg};
        border: 1px solid {border_color};
        border-radius: 16px;
        padding: 1.5rem;
        box-shadow: 0 4px 20px rgba(0, 0, 0, 0.2);
    }}

    [data-testid="stMetric"] label {{
        color: {text_muted} !important;
        font-size: 0.85rem !important;
        text-transform: uppercase;
        letter-spacing: 1px;
    }}

    [data-testid="stMetric"] [data-testid="stMetricValue"] {{
        color: {text_color} !important;
        font-size: 2rem !important;
        font-weight: 700 !important;
    }}

    [data-testid="stMetric"] [data-testid="stMetricDelta"] {{
        color: #22c55e !important;
    }}

    /* Tabs */
    .stTabs [data-baseweb="tab-list"] {{
        background: {card_bg};
        border: 1px solid {border_color};
        border-radius: 12px;
        padding: 0.5rem;
        gap: 0.5rem;
    }}

    .stTabs [data-baseweb="tab"] {{
        background: transparent;
        color: {text_muted};
        border-radius: 8px;
        padding: 0.75rem 1.5rem;
        font-weight: 500;
    }}

    .stTabs [aria-selected="true"] {{
        background: linear-gradient(135deg, #1e4976 0%, #2563eb 100%);
        color: white !important;
    }}

    /* Expander */
    .streamlit-expanderHeader {{
        background: {card_bg} !important;
        border: 1px solid {border_color};
        border-radius: 12px;
        color: {text_color} !important;
    }}

    .streamlit-expanderHeader p,
    .streamlit-expanderHeader span,
    [data-testid="stExpander"] summary,
    [data-testid="stExpander"] summary span,
    [data-testid="stExpander"] summary p {{
        color: {text_color} !important;
    }}

    [data-testid="stExpander"] {{
        background: {card_bg};
        border: 1px solid {border_color};
        border-radius: 12px;
    }}

    [data-testid="stExpander"] details {{
        background: {card_bg};
        border-radius: 12px;
    }}

    /* Dataframe */
    .stDataFrame {{
        background: {card_bg};
        border-radius: 12px;
    }}

    /* Divider */
    hr {{
        border-color: {border_color};
    }}

    /* Caption */
    .stCaption {{
        color: {text_muted} !important;
    }}

    /* Button */
    .stButton > button {{
        background: linear-gradient(135deg, #1e4976 0%, #2563eb 100%);
        color: white;
        border: none;
        border-radius: 12px;
        padding: 0.75rem 2rem;
        font-weight: 600;
        transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(37, 99, 235, 0.3);
    }}

    .stButton > button:hover {{
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(37, 99, 235, 0.4);
    }}

    /* Info box */
    .stAlert {{
        background: rgba(37, 99, 235, 0.1);
        border: 1px solid rgba(37, 99, 235, 0.25);
        border-radius: 12px;
    }}

    /* Plotly charts */
    .js-plotly-plot {{
        border-radius: 16px;
    }}

    /* Card styling */
    .premium-card {{
        background: {card_bg};
        border: 1px solid {border_color};
        border-radius: 16px;
        padding: 2rem;
        box-shadow: 0 4px 20px rgba(0, 0, 0, 0.2);
    }}

    /* ========== MOBILE RESPONSIVE ========== */
    @media (max-width: 768px) {{
        .main .block-container {{
            padding: 1rem;
        }}

        h1 {{
            font-size: 1.5rem !important;
        }}

        h2, h3 {{
            font-size: 1.2rem !important;
        }}

        [data-testid="stMetric"] {{
            padding: 0.75rem;
        }}

        [data-testid="stMetric"] [data-testid="stMetricValue"] {{
            font-size: 1.5rem !important;
        }}

        .stTabs [data-baseweb="tab"] {{
            padding: 0.5rem 0.75rem;
            font-size: 0.8rem;
        }}

        /* Stack columns on mobile */
        [data-testid="column"] {{
            width: 100% !important;
            flex: 1 1 100% !important;
        }}
    }}

    /* Sidebar mobile */
    @media (max-width: 992px) {{
        [data-testid="stSidebar"] {{
            width: 250px !important;
        }}
    }}
</style>
""", unsafe_allow_html=True)

# Layout do tema para gráficos Plotly - Navy Blue
PLOTLY_LAYOUT = dict(
    paper_bgcolor='rgba(0,0,0,0)',
    plot_bgcolor='rgba(0,0,0,0)',
    font=dict(color='#94a3b8', family='Inter, sans-serif'),
    title=dict(font=dict(color='#ffffff', size=18)),
    xaxis=dict(
        gridcolor='rgba(59, 130, 246, 0.1)',
        linecolor='rgba(59, 130, 246, 0.15)',
        tickfont=dict(color='#94a3b8')
    ),
    yaxis=dict(
        gridcolor='rgba(59, 130, 246, 0.1)',
        linecolor='rgba(59, 130, 246, 0.15)',
        tickfont=dict(color='#94a3b8')
    ),
    legend=dict(
        bgcolor='rgba(0,0,0,0)',
        font=dict(color='#94a3b8')
    ),
    margin=dict(t=60, b=40, l=40, r=40)
)

# Cores Corporate SaaS - Navy Blue
COLORS = {
    'primary': '#2563eb',      # Azul principal
    'secondary': '#1e4976',    # Azul marinho escuro
    'accent': '#3b82f6',       # Azul claro
    'success': '#22c55e',
    'warning': '#f59e0b',
    'danger': '#ef4444',
    'info': '#0ea5e9',
    'muted': '#64748b',
    'gradient': ['#1e4976', '#2563eb', '#3b82f6', '#60a5fa'],
    # Termômetro de ocupação (quanto maior, melhor)
    'hot': '#22c55e',      # 90-100% - Verde intenso (excelente)
    'warm': '#84cc16',     # 80-89% - Verde claro (muito bom)
    'mild': '#fbbf24',     # 70-79% - Amarelo (bom)
    'cool': '#f97316',     # 60-69% - Laranja (atenção)
    'cold': '#ef4444',     # <60% - Vermelho (crítico)
}

def get_ocupacao_color(ocupacao):
    """Retorna cor baseada na ocupação - quanto maior, mais quente/verde"""
    if ocupacao >= 90:
        return COLORS['hot']       # Verde intenso
    elif ocupacao >= 80:
        return COLORS['warm']      # Verde claro
    elif ocupacao >= 70:
        return COLORS['mild']      # Amarelo
    elif ocupacao >= 50:
        return COLORS['cool']      # Laranja
    else:
        return COLORS['cold']      # Vermelho (crítico)

BASE_PATH = Path(__file__).parent / "output"

# Carrega dados atuais
@st.cache_data(ttl=60)
def carregar_dados():
    with open(BASE_PATH / "resumo_ultimo.json") as f:
        resumo = json.load(f)
    with open(BASE_PATH / "vagas_ultimo.json") as f:
        vagas = json.load(f)
    return resumo, vagas

# Carrega histórico do banco
@st.cache_data(ttl=60)
def carregar_historico():
    db_path = BASE_PATH / "vagas.db"
    if not db_path.exists():
        return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), 0

    conn = sqlite3.connect(db_path)

    query_unidades = """
    SELECT e.data_extracao, v.unidade_codigo, v.unidade_nome,
           SUM(v.vagas) as vagas, SUM(v.matriculados) as matriculados,
           SUM(v.novatos) as novatos, SUM(v.veteranos) as veteranos,
           SUM(v.disponiveis) as disponiveis
    FROM vagas v JOIN 'extrações' e ON v.extracao_id = e.id
    GROUP BY e.id, v.unidade_codigo ORDER BY e.data_extracao
    """
    df_unidades = pd.read_sql_query(query_unidades, conn)

    query_total = """
    SELECT e.data_extracao, SUM(v.vagas) as vagas, SUM(v.matriculados) as matriculados,
           SUM(v.novatos) as novatos, SUM(v.veteranos) as veteranos, SUM(v.disponiveis) as disponiveis
    FROM vagas v JOIN 'extrações' e ON v.extracao_id = e.id
    GROUP BY e.id ORDER BY e.data_extracao
    """
    df_total = pd.read_sql_query(query_total, conn)

    query_segmento = """
    SELECT e.data_extracao, v.segmento, SUM(v.vagas) as vagas,
           SUM(v.matriculados) as matriculados, SUM(v.disponiveis) as disponiveis
    FROM vagas v JOIN 'extrações' e ON v.extracao_id = e.id
    GROUP BY e.id, v.segmento ORDER BY e.data_extracao
    """
    df_segmento = pd.read_sql_query(query_segmento, conn)

    cursor = conn.cursor()
    cursor.execute("SELECT COUNT(*) FROM 'extrações'")
    num_extracoes = cursor.fetchone()[0]
    conn.close()

    for df in [df_unidades, df_total, df_segmento]:
        if not df.empty:
            df['data_extracao'] = pd.to_datetime(df['data_extracao'])
            df['data_formatada'] = df['data_extracao'].dt.strftime('%d/%m %H:%M')

    return df_unidades, df_total, df_segmento, num_extracoes

try:
    resumo, vagas = carregar_dados()
    df_hist_unidades, df_hist_total, df_hist_segmento, num_extracoes = carregar_historico()
except FileNotFoundError:
    st.error("Arquivos de dados não encontrados. Execute a extração primeiro.")
    st.stop()

# Header Premium
col_title, col_btn = st.columns([5, 1])

with col_title:
    st.markdown("""
        <h1 style='margin-bottom: 0;'>Dashboard de Vagas</h1>
        <p style='color: #3b82f6; font-size: 1.2rem; margin-top: 0.5rem;'>Colégio Elo • Visão Executiva</p>
    """, unsafe_allow_html=True)

with col_btn:
    st.write("")
    if st.button("🔄 Atualizar", use_container_width=True):
        with st.spinner("Extraindo dados..."):
            try:
                result = subprocess.run(
                    ["bash", str(BASE_DIR / "cron_extrator.sh")],
                    capture_output=True, text=True, timeout=600, cwd=str(BASE_DIR)
                )
                if result.returncode == 0:
                    st.success("✅ Dados atualizados!")
                    st.cache_data.clear()
                    st.rerun()
                else:
                    st.error(f"❌ Erro: {result.stderr}")
            except Exception as e:
                st.error(f"❌ Erro: {str(e)}")

# Info bar
st.markdown(f"""
    <div style='display: flex; gap: 2rem; color: #64748b; font-size: 0.85rem; margin-bottom: 1rem;'>
        <span>📅 Última atualização: <strong style='color: #94a3b8;'>{resumo['data_extracao'][:16].replace('T', ' ')}</strong></span>
        <span>📊 Período: <strong style='color: #94a3b8;'>{resumo['periodo']}</strong></span>
        <span>🔢 Extrações: <strong style='color: #94a3b8;'>{num_extracoes}</strong></span>
    </div>
""", unsafe_allow_html=True)

# Botões de Download
col_dl1, col_dl2, col_dl3 = st.columns([1, 1, 4])

# Gerar Excel
def gerar_excel():
    output = BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # Resumo Geral
        df_resumo = pd.DataFrame([{
            'Métrica': 'Ocupação',
            'Valor': f"{round(resumo['total_geral']['matriculados'] / resumo['total_geral']['vagas'] * 100, 1)}%"
        }, {
            'Métrica': 'Matriculados',
            'Valor': resumo['total_geral']['matriculados']
        }, {
            'Métrica': 'Vagas Totais',
            'Valor': resumo['total_geral']['vagas']
        }, {
            'Métrica': 'Disponíveis',
            'Valor': resumo['total_geral']['vagas'] - resumo['total_geral']['matriculados']
        }, {
            'Métrica': 'Novatos',
            'Valor': resumo['total_geral']['novatos']
        }, {
            'Métrica': 'Veteranos',
            'Valor': resumo['total_geral']['veteranos']
        }])
        df_resumo.to_excel(writer, sheet_name='Resumo Geral', index=False)

        # Por Unidade
        dados_unidades = []
        for unidade in resumo['unidades']:
            nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
            t = unidade['total']
            ocup = round(t['matriculados'] / t['vagas'] * 100, 1)
            dados_unidades.append({
                'Unidade': nome,
                'Vagas': t['vagas'],
                'Novatos': t['novatos'],
                'Veteranos': t['veteranos'],
                'Matriculados': t['matriculados'],
                'Disponíveis': t['vagas'] - t['matriculados'],
                'Ocupação %': ocup
            })
        pd.DataFrame(dados_unidades).to_excel(writer, sheet_name='Por Unidade', index=False)

        # Todas as Turmas
        todas_turmas = []
        for unidade in vagas['unidades']:
            nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
            for turma in unidade['turmas']:
                ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0
                todas_turmas.append({
                    'Unidade': nome_unidade,
                    'Segmento': turma['segmento'],
                    'Turma': turma['turma'],
                    'Vagas': turma['vagas'],
                    'Novatos': turma['novatos'],
                    'Veteranos': turma['veteranos'],
                    'Matriculados': turma['matriculados'],
                    'Disponíveis': turma['vagas'] - turma['matriculados'],
                    'Pré-Matr.': turma['pre_matriculados'],
                    'Ocupação %': ocup
                })
        pd.DataFrame(todas_turmas).to_excel(writer, sheet_name='Todas as Turmas', index=False)

    return output.getvalue()

# Gerar PDF (HTML para impressão)
def gerar_pdf_html():
    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>Relatório de Vagas - Colégio Elo</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 40px; color: #333; }}
            h1 {{ color: #1e4976; border-bottom: 2px solid #2563eb; padding-bottom: 10px; }}
            h2 {{ color: #2563eb; margin-top: 30px; }}
            table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
            th, td {{ border: 1px solid #ddd; padding: 10px; text-align: center; }}
            th {{ background-color: #1e4976; color: white; }}
            tr:nth-child(even) {{ background-color: #f9f9f9; }}
            .metric {{ display: inline-block; margin: 10px 20px; text-align: center; }}
            .metric-value {{ font-size: 28px; font-weight: bold; color: #2563eb; }}
            .metric-label {{ font-size: 12px; color: #666; text-transform: uppercase; }}
            .footer {{ margin-top: 40px; text-align: center; color: #888; font-size: 12px; }}
            @media print {{ body {{ margin: 20px; }} }}
        </style>
    </head>
    <body>
        <h1>📊 Relatório de Vagas - Colégio Elo</h1>
        <p><strong>Data:</strong> {resumo['data_extracao'][:16].replace('T', ' ')} | <strong>Período:</strong> {resumo['periodo']}</p>

        <div style="background: #f0f4f8; padding: 20px; border-radius: 10px; margin: 20px 0;">
            <div class="metric">
                <div class="metric-value">{round(resumo['total_geral']['matriculados'] / resumo['total_geral']['vagas'] * 100, 1)}%</div>
                <div class="metric-label">Ocupação</div>
            </div>
            <div class="metric">
                <div class="metric-value">{resumo['total_geral']['matriculados']}</div>
                <div class="metric-label">Matriculados</div>
            </div>
            <div class="metric">
                <div class="metric-value">{resumo['total_geral']['vagas']}</div>
                <div class="metric-label">Vagas</div>
            </div>
            <div class="metric">
                <div class="metric-value">{resumo['total_geral']['vagas'] - resumo['total_geral']['matriculados']}</div>
                <div class="metric-label">Disponíveis</div>
            </div>
            <div class="metric">
                <div class="metric-value">{resumo['total_geral']['novatos']}</div>
                <div class="metric-label">Novatos</div>
            </div>
            <div class="metric">
                <div class="metric-value">{resumo['total_geral']['veteranos']}</div>
                <div class="metric-label">Veteranos</div>
            </div>
        </div>

        <h2>Por Unidade</h2>
        <table>
            <tr>
                <th>Unidade</th>
                <th>Vagas</th>
                <th>Matriculados</th>
                <th>Disponíveis</th>
                <th>Ocupação</th>
            </tr>
    """

    for unidade in resumo['unidades']:
        nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        t = unidade['total']
        ocup = round(t['matriculados'] / t['vagas'] * 100, 1)
        html += f"""
            <tr>
                <td><strong>{nome}</strong></td>
                <td>{t['vagas']}</td>
                <td>{t['matriculados']}</td>
                <td>{t['vagas'] - t['matriculados']}</td>
                <td>{ocup}%</td>
            </tr>
        """

    html += """
        </table>

        <h2>Por Segmento (Todas as Unidades)</h2>
        <table>
            <tr>
                <th>Segmento</th>
                <th>Vagas</th>
                <th>Matriculados</th>
                <th>Disponíveis</th>
                <th>Ocupação</th>
            </tr>
    """

    # Agrupa por segmento
    segmentos_totais = {}
    for unidade in resumo['unidades']:
        for seg, vals in unidade['segmentos'].items():
            if seg not in segmentos_totais:
                segmentos_totais[seg] = {'vagas': 0, 'matriculados': 0}
            segmentos_totais[seg]['vagas'] += vals['vagas']
            segmentos_totais[seg]['matriculados'] += vals['matriculados']

    for seg in ['Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']:
        if seg in segmentos_totais:
            v = segmentos_totais[seg]
            ocup = round(v['matriculados'] / v['vagas'] * 100, 1) if v['vagas'] > 0 else 0
            html += f"""
                <tr>
                    <td><strong>{seg}</strong></td>
                    <td>{v['vagas']}</td>
                    <td>{v['matriculados']}</td>
                    <td>{v['vagas'] - v['matriculados']}</td>
                    <td>{ocup}%</td>
                </tr>
            """

    html += f"""
        </table>

        <div class="footer">
            <p>Relatório gerado automaticamente pelo SIGA Vagas Dashboard</p>
            <p>Colégio Elo © {datetime.now().year}</p>
        </div>
    </body>
    </html>
    """
    return html

# ============================================================
# GERAÇÃO DE RELATÓRIOS PERSONALIZADOS
# ============================================================

def gerar_relatorio_resumo_executivo(formato='PDF'):
    """Gera relatório resumo executivo"""
    total = resumo['total_geral']
    ocupacao_geral = round(total['matriculados'] / total['vagas'] * 100, 1)

    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>Resumo Executivo - Colégio Elo</title>
        <style>
            @page {{ size: A4; margin: 2cm; }}
            body {{ font-family: 'Segoe UI', Arial, sans-serif; margin: 0; padding: 40px; color: #1e293b; background: #fff; }}
            .header {{ background: linear-gradient(135deg, #1e4976 0%, #2563eb 100%); color: white; padding: 30px; border-radius: 12px; margin-bottom: 30px; }}
            .header h1 {{ margin: 0; font-size: 28px; }}
            .header p {{ margin: 5px 0 0 0; opacity: 0.9; }}
            .kpi-grid {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 20px; margin: 30px 0; }}
            .kpi {{ background: #f8fafc; border: 1px solid #e2e8f0; border-radius: 12px; padding: 20px; text-align: center; }}
            .kpi-value {{ font-size: 36px; font-weight: 700; color: #2563eb; }}
            .kpi-label {{ font-size: 12px; color: #64748b; text-transform: uppercase; letter-spacing: 1px; margin-top: 8px; }}
            .section {{ margin: 30px 0; }}
            .section h2 {{ color: #1e4976; border-bottom: 2px solid #2563eb; padding-bottom: 10px; }}
            table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
            th {{ background: #1e4976; color: white; padding: 12px; text-align: center; }}
            td {{ border: 1px solid #e2e8f0; padding: 12px; text-align: center; }}
            tr:nth-child(even) {{ background: #f8fafc; }}
            .highlight {{ background: linear-gradient(90deg, #22c55e 0%, #84cc16 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent; font-weight: bold; }}
            .footer {{ text-align: center; margin-top: 40px; color: #94a3b8; font-size: 12px; }}
        </style>
    </head>
    <body>
        <div class="header">
            <h1>📊 Resumo Executivo</h1>
            <p>Colégio Elo • {datetime.now().strftime('%d/%m/%Y às %H:%M')}</p>
        </div>

        <div class="kpi-grid">
            <div class="kpi">
                <div class="kpi-value">{ocupacao_geral}%</div>
                <div class="kpi-label">Ocupação Geral</div>
            </div>
            <div class="kpi">
                <div class="kpi-value">{total['matriculados']:,}</div>
                <div class="kpi-label">Matriculados</div>
            </div>
            <div class="kpi">
                <div class="kpi-value">{total['disponiveis']:,}</div>
                <div class="kpi-label">Disponíveis</div>
            </div>
        </div>

        <div class="section">
            <h2>Desempenho por Unidade</h2>
            <table>
                <tr>
                    <th>Unidade</th>
                    <th>Vagas</th>
                    <th>Matriculados</th>
                    <th>Disponíveis</th>
                    <th>Ocupação</th>
                </tr>
    """

    for unidade in resumo['unidades']:
        nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        t = unidade['total']
        ocup = round(t['matriculados'] / t['vagas'] * 100, 1)
        html += f"""
                <tr>
                    <td><strong>{nome}</strong></td>
                    <td>{t['vagas']}</td>
                    <td>{t['matriculados']}</td>
                    <td>{t['disponiveis']}</td>
                    <td><strong>{ocup}%</strong></td>
                </tr>
        """

    html += f"""
                <tr style="background: #1e4976; color: white;">
                    <td><strong>TOTAL</strong></td>
                    <td><strong>{total['vagas']}</strong></td>
                    <td><strong>{total['matriculados']}</strong></td>
                    <td><strong>{total['disponiveis']}</strong></td>
                    <td><strong>{ocupacao_geral}%</strong></td>
                </tr>
            </table>
        </div>

        <div class="footer">
            <p>Relatório gerado automaticamente pelo SIGA Vagas Dashboard</p>
            <p>Colégio Elo © {datetime.now().year}</p>
        </div>
    </body>
    </html>
    """
    return html

def gerar_relatorio_detalhado_unidade(formato='PDF'):
    """Gera relatório detalhado por unidade"""
    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>Relatório Detalhado por Unidade - Colégio Elo</title>
        <style>
            @page {{ size: A4; margin: 1.5cm; }}
            body {{ font-family: 'Segoe UI', Arial, sans-serif; margin: 0; padding: 30px; color: #1e293b; font-size: 11px; }}
            .header {{ background: linear-gradient(135deg, #1e4976 0%, #2563eb 100%); color: white; padding: 20px; border-radius: 8px; margin-bottom: 20px; }}
            .header h1 {{ margin: 0; font-size: 22px; }}
            .unidade-section {{ margin: 20px 0; page-break-inside: avoid; }}
            .unidade-header {{ background: #f1f5f9; padding: 12px; border-radius: 8px; margin-bottom: 10px; }}
            .unidade-header h3 {{ margin: 0; color: #1e4976; }}
            table {{ width: 100%; border-collapse: collapse; font-size: 10px; }}
            th {{ background: #334155; color: white; padding: 8px; text-align: center; }}
            td {{ border: 1px solid #e2e8f0; padding: 6px; text-align: center; }}
            .critico {{ background: #fef2f2; color: #dc2626; font-weight: bold; }}
            .atencao {{ background: #fffbeb; color: #d97706; }}
            .bom {{ background: #f0fdf4; color: #16a34a; }}
            .footer {{ text-align: center; margin-top: 30px; color: #94a3b8; font-size: 10px; }}
        </style>
    </head>
    <body>
        <div class="header">
            <h1>📋 Relatório Detalhado por Unidade</h1>
            <p>Colégio Elo • {datetime.now().strftime('%d/%m/%Y às %H:%M')}</p>
        </div>
    """

    for unidade in vagas['unidades']:
        nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        total_unid = next((u['total'] for u in resumo['unidades'] if u['codigo'] == unidade['codigo']), {})
        ocup_unid = round(total_unid.get('matriculados', 0) / total_unid.get('vagas', 1) * 100, 1)

        html += f"""
        <div class="unidade-section">
            <div class="unidade-header">
                <h3>🏫 {nome_unidade} - Ocupação: {ocup_unid}%</h3>
            </div>
            <table>
                <tr>
                    <th>Turma</th>
                    <th>Segmento</th>
                    <th>Vagas</th>
                    <th>Matr.</th>
                    <th>Disp.</th>
                    <th>Ocup.</th>
                </tr>
        """

        for turma in unidade['turmas']:
            ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0
            classe = 'critico' if ocup < 50 else ('atencao' if ocup < 70 else ('bom' if ocup >= 80 else ''))
            html += f"""
                <tr class="{classe}">
                    <td style="text-align: left;">{turma['turma'][:40]}</td>
                    <td>{turma['segmento']}</td>
                    <td>{turma['vagas']}</td>
                    <td>{turma['matriculados']}</td>
                    <td>{turma['vagas'] - turma['matriculados']}</td>
                    <td><strong>{ocup}%</strong></td>
                </tr>
            """

        html += "</table></div>"

    html += f"""
        <div class="footer">
            <p>Relatório gerado automaticamente pelo SIGA Vagas Dashboard • Colégio Elo © {datetime.now().year}</p>
        </div>
    </body>
    </html>
    """
    return html

def gerar_relatorio_turmas_criticas(formato='PDF'):
    """Gera relatório de turmas críticas"""
    # Coleta turmas críticas
    turmas_report = []
    for unidade in vagas['unidades']:
        nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        for turma in unidade['turmas']:
            ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0
            if ocup < 70:  # Turmas com ocupação abaixo de 70%
                turmas_report.append({
                    'unidade': nome_unidade,
                    'turma': turma['turma'],
                    'segmento': turma['segmento'],
                    'vagas': turma['vagas'],
                    'matriculados': turma['matriculados'],
                    'disponiveis': turma['vagas'] - turma['matriculados'],
                    'ocupacao': ocup
                })

    turmas_report = sorted(turmas_report, key=lambda x: x['ocupacao'])

    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>Turmas Críticas - Colégio Elo</title>
        <style>
            body {{ font-family: 'Segoe UI', Arial, sans-serif; margin: 0; padding: 40px; color: #1e293b; }}
            .header {{ background: linear-gradient(135deg, #dc2626 0%, #f97316 100%); color: white; padding: 30px; border-radius: 12px; margin-bottom: 30px; }}
            .header h1 {{ margin: 0; font-size: 28px; }}
            .alerta {{ background: #fef2f2; border: 2px solid #dc2626; border-radius: 12px; padding: 20px; margin-bottom: 20px; }}
            .alerta h2 {{ color: #dc2626; margin: 0 0 10px 0; }}
            table {{ width: 100%; border-collapse: collapse; }}
            th {{ background: #dc2626; color: white; padding: 12px; text-align: center; }}
            td {{ border: 1px solid #fecaca; padding: 10px; text-align: center; }}
            tr:nth-child(even) {{ background: #fef2f2; }}
            .ocupacao-baixa {{ color: #dc2626; font-weight: bold; font-size: 16px; }}
            .footer {{ text-align: center; margin-top: 40px; color: #94a3b8; font-size: 12px; }}
        </style>
    </head>
    <body>
        <div class="header">
            <h1>🚨 Relatório de Turmas Críticas</h1>
            <p>Colégio Elo • {datetime.now().strftime('%d/%m/%Y às %H:%M')}</p>
        </div>

        <div class="alerta">
            <h2>⚠️ {len(turmas_report)} turmas com ocupação abaixo de 70%</h2>
            <p>Estas turmas requerem atenção especial para atingir as metas de matrícula.</p>
        </div>

        <table>
            <tr>
                <th>Unidade</th>
                <th>Turma</th>
                <th>Segmento</th>
                <th>Vagas</th>
                <th>Matr.</th>
                <th>Disp.</th>
                <th>Ocupação</th>
            </tr>
    """

    for t in turmas_report[:30]:  # Limita a 30 turmas
        html += f"""
            <tr>
                <td>{t['unidade']}</td>
                <td style="text-align: left;">{t['turma'][:35]}</td>
                <td>{t['segmento']}</td>
                <td>{t['vagas']}</td>
                <td>{t['matriculados']}</td>
                <td>{t['disponiveis']}</td>
                <td class="ocupacao-baixa">{t['ocupacao']}%</td>
            </tr>
        """

    html += f"""
        </table>

        <div class="footer">
            <p>Relatório gerado automaticamente pelo SIGA Vagas Dashboard</p>
            <p>Colégio Elo © {datetime.now().year}</p>
        </div>
    </body>
    </html>
    """
    return html

def gerar_relatorio_tendencias(formato='PDF'):
    """Gera relatório de análise de tendências"""
    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>Análise de Tendências - Colégio Elo</title>
        <style>
            body {{ font-family: 'Segoe UI', Arial, sans-serif; margin: 0; padding: 40px; color: #1e293b; }}
            .header {{ background: linear-gradient(135deg, #059669 0%, #10b981 100%); color: white; padding: 30px; border-radius: 12px; margin-bottom: 30px; }}
            .header h1 {{ margin: 0; font-size: 28px; }}
            .section {{ margin: 30px 0; padding: 20px; background: #f8fafc; border-radius: 12px; }}
            .section h2 {{ color: #1e4976; margin: 0 0 15px 0; }}
            table {{ width: 100%; border-collapse: collapse; }}
            th {{ background: #334155; color: white; padding: 12px; text-align: center; }}
            td {{ border: 1px solid #e2e8f0; padding: 10px; text-align: center; }}
            .trend-up {{ color: #16a34a; font-weight: bold; }}
            .trend-down {{ color: #dc2626; font-weight: bold; }}
            .footer {{ text-align: center; margin-top: 40px; color: #94a3b8; font-size: 12px; }}
        </style>
    </head>
    <body>
        <div class="header">
            <h1>📈 Análise de Tendências</h1>
            <p>Colégio Elo • {datetime.now().strftime('%d/%m/%Y às %H:%M')}</p>
        </div>

        <div class="section">
            <h2>📊 Ocupação por Segmento</h2>
            <table>
                <tr>
                    <th>Segmento</th>
                    <th>Vagas</th>
                    <th>Matriculados</th>
                    <th>Disponíveis</th>
                    <th>Ocupação</th>
                </tr>
    """

    # Agrupa por segmento
    segmentos_totais = {}
    for unidade in resumo['unidades']:
        for seg, vals in unidade['segmentos'].items():
            if seg not in segmentos_totais:
                segmentos_totais[seg] = {'vagas': 0, 'matriculados': 0}
            segmentos_totais[seg]['vagas'] += vals['vagas']
            segmentos_totais[seg]['matriculados'] += vals['matriculados']

    for seg in ['Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']:
        if seg in segmentos_totais:
            v = segmentos_totais[seg]
            ocup = round(v['matriculados'] / v['vagas'] * 100, 1) if v['vagas'] > 0 else 0
            classe = 'trend-up' if ocup >= 70 else 'trend-down'
            html += f"""
                <tr>
                    <td><strong>{seg}</strong></td>
                    <td>{v['vagas']}</td>
                    <td>{v['matriculados']}</td>
                    <td>{v['vagas'] - v['matriculados']}</td>
                    <td class="{classe}">{ocup}%</td>
                </tr>
            """

    total = resumo['total_geral']
    html += f"""
            </table>
        </div>

        <div class="section">
            <h2>🎯 Projeções</h2>
            <p><strong>Ocupação atual:</strong> {round(total['matriculados'] / total['vagas'] * 100, 1)}%</p>
            <p><strong>Meta de ocupação:</strong> 80%</p>
            <p><strong>Matrículas necessárias para meta:</strong> {max(0, int(total['vagas'] * 0.8) - total['matriculados'])}</p>
            <p><strong>Vagas disponíveis:</strong> {total['disponiveis']}</p>
        </div>

        <div class="footer">
            <p>Relatório gerado automaticamente pelo SIGA Vagas Dashboard</p>
            <p>Colégio Elo © {datetime.now().year}</p>
        </div>
    </body>
    </html>
    """
    return html

def gerar_excel_relatorio(tipo_relatorio):
    """Gera relatório em formato Excel"""
    output = BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # Resumo Geral sempre incluído
        total = resumo['total_geral']
        df_resumo = pd.DataFrame([{
            'Métrica': 'Ocupação',
            'Valor': f"{round(total['matriculados'] / total['vagas'] * 100, 1)}%"
        }, {
            'Métrica': 'Matriculados',
            'Valor': total['matriculados']
        }, {
            'Métrica': 'Vagas Totais',
            'Valor': total['vagas']
        }, {
            'Métrica': 'Disponíveis',
            'Valor': total['disponiveis']
        }])
        df_resumo.to_excel(writer, sheet_name='Resumo', index=False)

        # Por Unidade
        dados_unidades = []
        for unidade in resumo['unidades']:
            nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
            t = unidade['total']
            ocup = round(t['matriculados'] / t['vagas'] * 100, 1)
            dados_unidades.append({
                'Unidade': nome,
                'Vagas': t['vagas'],
                'Matriculados': t['matriculados'],
                'Disponíveis': t['disponiveis'],
                'Ocupação %': ocup
            })
        pd.DataFrame(dados_unidades).to_excel(writer, sheet_name='Por Unidade', index=False)

        if tipo_relatorio in ['Detalhado por Unidade', 'Turmas Críticas']:
            # Todas as turmas
            todas_turmas = []
            for unidade in vagas['unidades']:
                nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
                for turma in unidade['turmas']:
                    ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0
                    if tipo_relatorio == 'Turmas Críticas' and ocup >= 70:
                        continue
                    todas_turmas.append({
                        'Unidade': nome_unidade,
                        'Segmento': turma['segmento'],
                        'Turma': turma['turma'],
                        'Vagas': turma['vagas'],
                        'Matriculados': turma['matriculados'],
                        'Disponíveis': turma['vagas'] - turma['matriculados'],
                        'Ocupação %': ocup
                    })
            if todas_turmas:
                df_turmas = pd.DataFrame(todas_turmas)
                df_turmas = df_turmas.sort_values('Ocupação %', ascending=True)
                df_turmas.to_excel(writer, sheet_name='Turmas', index=False)

    return output.getvalue()

def gerar_powerpoint(tipo_relatorio):
    """Gera relatório em formato PowerPoint"""
    if not PPTX_AVAILABLE:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Cores do tema
    AZUL_ESCURO = RGBColor(30, 73, 118)  # #1e4976
    AZUL_CLARO = RGBColor(37, 99, 235)   # #2563eb
    VERDE = RGBColor(34, 197, 94)        # #22c55e
    VERMELHO = RGBColor(239, 68, 68)     # #ef4444
    LARANJA = RGBColor(249, 115, 22)     # #f97316
    CINZA = RGBColor(100, 116, 139)      # #64748b
    BRANCO = RGBColor(255, 255, 255)

    total = resumo['total_geral']
    ocupacao_geral = round(total['matriculados'] / total['vagas'] * 100, 1)

    # ========== SLIDE 1: CAPA ==========
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZUL_ESCURO
    bg.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📊 {tipo_relatorio}"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Colégio Elo • Sistema de Gestão de Vagas"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    # Data
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime('%d/%m/%Y às %H:%M')
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 2: KPIs ==========
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_ESCURO
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📈 Indicadores Principais"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # KPIs
    kpis = [
        ("Ocupação", f"{ocupacao_geral}%", VERDE if ocupacao_geral >= 70 else VERMELHO),
        ("Matriculados", f"{total['matriculados']:,}".replace(",", "."), AZUL_CLARO),
        ("Vagas Totais", f"{total['vagas']:,}".replace(",", "."), AZUL_CLARO),
        ("Disponíveis", f"{total['disponiveis']:,}".replace(",", "."), LARANJA),
    ]

    kpi_width = Inches(2.8)
    kpi_height = Inches(2)
    start_x = Inches(0.8)
    start_y = Inches(2)
    gap = Inches(0.4)

    for i, (label, value, color) in enumerate(kpis):
        x = start_x + i * (kpi_width + gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, kpi_width, kpi_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(241, 245, 249)
        card.line.color.rgb = RGBColor(226, 232, 240)

        # Value
        val_box = slide.shapes.add_textbox(x, start_y + Inches(0.3), kpi_width, Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x, start_y + Inches(1.3), kpi_width, Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA
        p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 3: POR UNIDADE ==========
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_ESCURO
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "🏫 Desempenho por Unidade"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # Tabela de unidades
    rows = len(resumo['unidades']) + 2  # Header + unidades + total
    cols = 5
    table_width = Inches(11)
    table_height = Inches(0.5) * rows

    table = slide.shapes.add_table(rows, cols, Inches(1.1), Inches(1.8), table_width, table_height).table

    # Header da tabela
    headers = ['Unidade', 'Vagas', 'Matriculados', 'Disponíveis', 'Ocupação']
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_ESCURO
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

    # Dados das unidades
    for i, unidade in enumerate(resumo['unidades'], 1):
        nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        t = unidade['total']
        ocup = round(t['matriculados'] / t['vagas'] * 100, 1)

        dados = [nome, str(t['vagas']), str(t['matriculados']), str(t['disponiveis']), f"{ocup}%"]

        for j, d in enumerate(dados):
            cell = table.cell(i, j)
            cell.text = d
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

            # Cor da ocupação
            if j == 4:
                if ocup >= 80:
                    p.font.color.rgb = VERDE
                elif ocup >= 50:
                    p.font.color.rgb = LARANJA
                else:
                    p.font.color.rgb = VERMELHO
                p.font.bold = True

    # Linha total
    total_row = rows - 1
    totais = ['TOTAL', str(total['vagas']), str(total['matriculados']), str(total['disponiveis']), f"{ocupacao_geral}%"]
    for j, d in enumerate(totais):
        cell = table.cell(total_row, j)
        cell.text = d
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

    # ========== SLIDE 4: POR SEGMENTO ==========
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_ESCURO
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📚 Desempenho por Segmento"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # Agrupa por segmento
    segmentos_totais = {}
    for unidade in resumo['unidades']:
        for seg, vals in unidade['segmentos'].items():
            if seg not in segmentos_totais:
                segmentos_totais[seg] = {'vagas': 0, 'matriculados': 0}
            segmentos_totais[seg]['vagas'] += vals['vagas']
            segmentos_totais[seg]['matriculados'] += vals['matriculados']

    # Cards de segmentos
    seg_ordem = ['Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']
    card_width = Inches(2.8)
    card_height = Inches(3.5)
    start_x = Inches(0.8)
    gap = Inches(0.4)

    for i, seg in enumerate(seg_ordem):
        if seg not in segmentos_totais:
            continue

        v = segmentos_totais[seg]
        ocup = round(v['matriculados'] / v['vagas'] * 100, 1) if v['vagas'] > 0 else 0
        x = start_x + i * (card_width + gap)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        card.line.color.rgb = RGBColor(226, 232, 240)

        # Título do segmento
        seg_title = slide.shapes.add_textbox(x, Inches(2.2), card_width, Inches(0.5))
        tf = seg_title.text_frame
        p = tf.paragraphs[0]
        p.text = seg
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = AZUL_ESCURO
        p.alignment = PP_ALIGN.CENTER

        # Ocupação
        ocup_box = slide.shapes.add_textbox(x, Inches(2.8), card_width, Inches(0.8))
        tf = ocup_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{ocup}%"
        p.font.size = Pt(36)
        p.font.bold = True
        if ocup >= 80:
            p.font.color.rgb = VERDE
        elif ocup >= 50:
            p.font.color.rgb = LARANJA
        else:
            p.font.color.rgb = VERMELHO
        p.alignment = PP_ALIGN.CENTER

        # Detalhes
        det_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.8), card_width - Inches(0.4), Inches(1.5))
        tf = det_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"Vagas: {v['vagas']}"
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA

        p = tf.add_paragraph()
        p.text = f"Matriculados: {v['matriculados']}"
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA

        p = tf.add_paragraph()
        p.text = f"Disponíveis: {v['vagas'] - v['matriculados']}"
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA

    # ========== SLIDE 5: TURMAS CRÍTICAS (se aplicável) ==========
    if tipo_relatorio in ['Turmas Críticas', 'Detalhado por Unidade']:
        slide = prs.slides.add_slide(slide_layout)

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = VERMELHO
        header.line.fill.background()

        header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = header_text.text_frame
        p = tf.paragraphs[0]
        p.text = "🚨 Turmas que Requerem Atenção"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BRANCO

        # Coleta turmas críticas
        turmas_criticas_ppt = []
        for unidade in vagas['unidades']:
            nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
            for turma in unidade['turmas']:
                ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0
                if ocup < 70:
                    turmas_criticas_ppt.append({
                        'unidade': nome_unidade,
                        'turma': turma['turma'][:35],
                        'segmento': turma['segmento'],
                        'ocupacao': ocup
                    })

        turmas_criticas_ppt = sorted(turmas_criticas_ppt, key=lambda x: x['ocupacao'])[:10]

        if turmas_criticas_ppt:
            rows = len(turmas_criticas_ppt) + 1
            table = slide.shapes.add_table(rows, 4, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.45) * rows).table

            headers = ['Unidade', 'Turma', 'Segmento', 'Ocupação']
            for j, h in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = h
                cell.fill.solid()
                cell.fill.fore_color.rgb = VERMELHO
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = BRANCO
                p.alignment = PP_ALIGN.CENTER

            for i, t in enumerate(turmas_criticas_ppt, 1):
                dados = [t['unidade'], t['turma'], t['segmento'], f"{t['ocupacao']}%"]
                for j, d in enumerate(dados):
                    cell = table.cell(i, j)
                    cell.text = d
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(11)
                    p.alignment = PP_ALIGN.CENTER if j > 1 else PP_ALIGN.LEFT
                    if j == 3:
                        p.font.bold = True
                        p.font.color.rgb = VERMELHO

    # ========== SLIDE 6: GRÁFICO DE PIZZA - OCUPAÇÃO ==========
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_ESCURO
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 Distribuição de Vagas"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # Gráfico de Pizza - Matriculados vs Disponíveis
    chart_data = CategoryChartData()
    chart_data.categories = ['Matriculados', 'Disponíveis']
    chart_data.add_series('Vagas', (total['matriculados'], total['disponiveis']))

    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(5.5), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Cores do gráfico de pizza
    series = chart.series[0]
    series.points[0].format.fill.solid()
    series.points[0].format.fill.fore_color.rgb = VERDE  # Matriculados
    series.points[1].format.fill.solid()
    series.points[1].format.fill.fore_color.rgb = LARANJA  # Disponíveis

    # Labels com porcentagem
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = False
    data_labels.font.size = Pt(14)
    data_labels.font.bold = True

    # KPI ao lado do gráfico
    kpi_box = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5), Inches(3))
    tf = kpi_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"{ocupacao_geral}%"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = VERDE if ocupacao_geral >= 70 else VERMELHO
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Taxa de Ocupação"
    p.font.size = Pt(18)
    p.font.color.rgb = CINZA
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = f"\n{total['matriculados']:,} matriculados".replace(",", ".")
    p.font.size = Pt(16)
    p.font.color.rgb = CINZA
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = f"{total['disponiveis']:,} disponíveis".replace(",", ".")
    p.font.size = Pt(16)
    p.font.color.rgb = CINZA
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 7: GRÁFICO DE BARRAS - POR UNIDADE ==========
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_ESCURO
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 Ocupação por Unidade"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # Dados do gráfico
    chart_data = CategoryChartData()
    nomes_unidades = []
    ocupacoes_unidades = []
    for unidade in resumo['unidades']:
        nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        t = unidade['total']
        ocup = round(t['matriculados'] / t['vagas'] * 100, 1)
        nomes_unidades.append(nome)
        ocupacoes_unidades.append(ocup)

    chart_data.categories = nomes_unidades
    chart_data.add_series('Ocupação %', ocupacoes_unidades)

    x, y, cx, cy = Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False

    # Configurar eixo Y
    value_axis = chart.value_axis
    value_axis.maximum_scale = 100
    value_axis.minimum_scale = 0
    value_axis.major_unit = 20
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(226, 232, 240)

    # Data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.font.size = Pt(14)
    data_labels.font.bold = True
    data_labels.number_format = '0"%"'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # Cores das barras
    series = chart.series[0]
    for i, ocup in enumerate(ocupacoes_unidades):
        point = series.points[i]
        point.format.fill.solid()
        if ocup >= 80:
            point.format.fill.fore_color.rgb = VERDE
        elif ocup >= 50:
            point.format.fill.fore_color.rgb = LARANJA
        else:
            point.format.fill.fore_color.rgb = VERMELHO

    # Linha de meta (80%)
    meta_box = slide.shapes.add_textbox(Inches(10.5), Inches(1.4), Inches(2), Inches(0.3))
    tf = meta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Meta: 80%"
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA

    # ========== SLIDE 8: GRÁFICO DE BARRAS - POR SEGMENTO ==========
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_ESCURO
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 Ocupação por Segmento"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # Dados do gráfico por segmento
    chart_data = CategoryChartData()
    seg_ordem = ['Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']
    nomes_seg = []
    ocupacoes_seg = []

    for seg in seg_ordem:
        if seg in segmentos_totais:
            v = segmentos_totais[seg]
            ocup = round(v['matriculados'] / v['vagas'] * 100, 1) if v['vagas'] > 0 else 0
            nomes_seg.append(seg)
            ocupacoes_seg.append(ocup)

    chart_data.categories = nomes_seg
    chart_data.add_series('Ocupação %', ocupacoes_seg)

    x, y, cx, cy = Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False

    # Configurar eixo
    value_axis = chart.value_axis
    value_axis.maximum_scale = 100
    value_axis.minimum_scale = 0
    value_axis.major_unit = 20
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(226, 232, 240)

    # Data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.font.size = Pt(14)
    data_labels.font.bold = True
    data_labels.number_format = '0"%"'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # Cores das barras por segmento
    series = chart.series[0]
    for i, ocup in enumerate(ocupacoes_seg):
        point = series.points[i]
        point.format.fill.solid()
        if ocup >= 80:
            point.format.fill.fore_color.rgb = VERDE
        elif ocup >= 50:
            point.format.fill.fore_color.rgb = LARANJA
        else:
            point.format.fill.fore_color.rgb = VERMELHO

    # ========== SLIDE FINAL ==========
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZUL_ESCURO
    bg.line.fill.background()

    # Obrigado
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Info
    info = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "Relatório gerado pelo SIGA Vagas Dashboard"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = f"Colégio Elo © {datetime.now().year}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    # Salva em BytesIO
    output = BytesIO()
    prs.save(output)
    return output.getvalue()

with col_dl1:
    excel_data = gerar_excel()
    st.download_button(
        label="📥 Excel",
        data=excel_data,
        file_name=f"vagas_colegio_elo_{datetime.now().strftime('%Y%m%d')}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True
    )

with col_dl2:
    pdf_html = gerar_pdf_html()
    st.download_button(
        label="📄 Relatório",
        data=pdf_html,
        file_name=f"relatorio_vagas_{datetime.now().strftime('%Y%m%d')}.html",
        mime="text/html",
        use_container_width=True,
        help="Abra o arquivo e use Ctrl+P para imprimir como PDF"
    )

# ============================================================
# HANDLER DE GERAÇÃO DE RELATÓRIOS PERSONALIZADOS
# ============================================================
if st.session_state.get('gerar_relatorio', False):
    tipo = st.session_state.get('tipo_relatorio_selecionado', 'Resumo Executivo')
    formato = st.session_state.get('formato_export_selecionado', 'PDF')

    st.markdown("---")
    st.markdown(f"### 📄 Relatório Gerado: {tipo}")

    # Gera o conteúdo baseado no tipo
    if tipo == "Resumo Executivo":
        html_content = gerar_relatorio_resumo_executivo(formato)
    elif tipo == "Detalhado por Unidade":
        html_content = gerar_relatorio_detalhado_unidade(formato)
    elif tipo == "Análise de Tendências":
        html_content = gerar_relatorio_tendencias(formato)
    else:  # Turmas Críticas
        html_content = gerar_relatorio_turmas_criticas(formato)

    # Oferece download baseado no formato
    col_rel1, col_rel2, col_rel3 = st.columns([2, 2, 2])

    with col_rel1:
        if formato == "PDF":
            st.download_button(
                label="📄 Baixar PDF (HTML)",
                data=html_content,
                file_name=f"relatorio_{tipo.lower().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.html",
                mime="text/html",
                use_container_width=True,
                help="Abra no navegador e use Ctrl+P para salvar como PDF"
            )
        elif formato == "Excel":
            excel_data = gerar_excel_relatorio(tipo)
            st.download_button(
                label="📊 Baixar Excel",
                data=excel_data,
                file_name=f"relatorio_{tipo.lower().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
        else:  # PowerPoint
            if PPTX_AVAILABLE:
                pptx_data = gerar_powerpoint(tipo)
                if pptx_data:
                    st.download_button(
                        label="📊 Baixar PowerPoint",
                        data=pptx_data,
                        file_name=f"relatorio_{tipo.lower().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                else:
                    st.error("Erro ao gerar PowerPoint")
            else:
                st.warning("⚠️ Biblioteca python-pptx não instalada. Use: `pip install python-pptx`")
                st.download_button(
                    label="📄 Baixar PDF (alternativa)",
                    data=html_content,
                    file_name=f"relatorio_{tipo.lower().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.html",
                    mime="text/html",
                    use_container_width=True
                )

    with col_rel2:
        if st.button("❌ Fechar", use_container_width=True, key="btn_fechar_relatorio"):
            st.session_state['gerar_relatorio'] = False
            st.rerun()

    # Preview do relatório
    with st.expander("👁️ Pré-visualização do Relatório", expanded=False):
        st.components.v1.html(html_content, height=600, scrolling=True)

    st.markdown("---")

st.markdown("<br>", unsafe_allow_html=True)

# Métricas principais
total = resumo['total_geral']
ocupacao = round(total['matriculados'] / total['vagas'] * 100, 1)

col1, col2, col3, col4, col5, col6 = st.columns(6)

with col1:
    st.metric("OCUPAÇÃO", f"{ocupacao}%", delta=None)
with col2:
    st.metric("MATRICULADOS", f"{total['matriculados']:,}".replace(",", "."))
with col3:
    st.metric("VAGAS TOTAIS", f"{total['vagas']:,}".replace(",", "."))
with col4:
    st.metric("DISPONÍVEIS", f"{total['disponiveis']:,}".replace(",", "."))
with col5:
    st.metric("NOVATOS", f"{total['novatos']:,}".replace(",", "."))
with col6:
    st.metric("VETERANOS", f"{total['veteranos']:,}".replace(",", "."))

st.markdown("<br>", unsafe_allow_html=True)

# Alertas de Turmas
todas_turmas_alerta = []
for unidade in vagas['unidades']:
    nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
    for turma in unidade['turmas']:
        ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0
        todas_turmas_alerta.append({
            'unidade': nome_unidade,
            'segmento': turma['segmento'],
            'turma': turma['turma'],
            'vagas': turma['vagas'],
            'matriculados': turma['matriculados'],
            'disponiveis': turma['vagas'] - turma['matriculados'],
            'ocupacao': ocup
        })

# Turmas críticas e quase lotadas (usando limites configuráveis)
turmas_criticas = sorted([t for t in todas_turmas_alerta if t['ocupacao'] < alerta_critico], key=lambda x: x['ocupacao'])[:5]
turmas_atencao = sorted([t for t in todas_turmas_alerta if alerta_critico <= t['ocupacao'] < alerta_atencao], key=lambda x: x['ocupacao'])[:5]
turmas_lotadas = sorted([t for t in todas_turmas_alerta if t['ocupacao'] >= alerta_lotado], key=lambda x: -x['ocupacao'])[:5]

# Contadores para histórico
total_criticas = len([t for t in todas_turmas_alerta if t['ocupacao'] < alerta_critico])
total_atencao = len([t for t in todas_turmas_alerta if alerta_critico <= t['ocupacao'] < alerta_atencao])
total_lotadas = len([t for t in todas_turmas_alerta if t['ocupacao'] >= alerta_lotado])

# Painel de Alertas
st.markdown(f"""
    <div style='background: linear-gradient(90deg, rgba(239, 68, 68, 0.1) 0%, rgba(251, 191, 36, 0.1) 50%, rgba(34, 197, 94, 0.1) 100%);
                padding: 1rem; border-radius: 12px; margin-bottom: 1rem;'>
        <div style='display: flex; justify-content: space-around; text-align: center;'>
            <div>
                <div style='font-size: 2rem; font-weight: bold; color: #ef4444;'>{total_criticas}</div>
                <div style='color: #94a3b8; font-size: 0.8rem;'>❄️ Críticas (&lt;{alerta_critico}%)</div>
            </div>
            <div>
                <div style='font-size: 2rem; font-weight: bold; color: #f97316;'>{total_atencao}</div>
                <div style='color: #94a3b8; font-size: 0.8rem;'>⚠️ Atenção ({alerta_critico}-{alerta_atencao}%)</div>
            </div>
            <div>
                <div style='font-size: 2rem; font-weight: bold; color: #22c55e;'>{total_lotadas}</div>
                <div style='color: #94a3b8; font-size: 0.8rem;'>🔥 Lotadas (&gt;{alerta_lotado}%)</div>
            </div>
        </div>
    </div>
""", unsafe_allow_html=True)

if turmas_criticas or turmas_lotadas:
    col_alert1, col_alert2 = st.columns(2)

    with col_alert1:
        if turmas_criticas:
            st.markdown(f"""
                <div style='background: linear-gradient(135deg, rgba(239, 68, 68, 0.15) 0%, rgba(239, 68, 68, 0.05) 100%);
                            border: 1px solid rgba(239, 68, 68, 0.3); border-radius: 12px; padding: 1rem;'>
                    <h4 style='color: #ef4444; margin: 0 0 0.5rem 0;'>❄️ Turmas Críticas (&lt;{alerta_critico}%)</h4>
            """, unsafe_allow_html=True)
            for t in turmas_criticas:
                st.markdown(f"""
                    <div style='background: rgba(0,0,0,0.2); border-radius: 8px; padding: 0.5rem; margin: 0.3rem 0;'>
                        <span style='color: #ef4444; font-weight: bold;'>{t['ocupacao']}%</span>
                        <span style='color: #94a3b8;'> • {t['unidade']} • {t['segmento']}</span><br>
                        <span style='color: #fff; font-size: 0.85rem;'>{t['turma'][:50]}...</span>
                    </div>
                """, unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)
        else:
            st.markdown("""
                <div style='background: rgba(34, 197, 94, 0.1); border: 1px solid rgba(34, 197, 94, 0.3);
                            border-radius: 12px; padding: 1rem; text-align: center;'>
                    <span style='color: #22c55e; font-size: 1.2rem;'>✅ Nenhuma turma crítica!</span>
                </div>
            """, unsafe_allow_html=True)

    with col_alert2:
        if turmas_lotadas:
            st.markdown(f"""
                <div style='background: linear-gradient(135deg, rgba(34, 197, 94, 0.15) 0%, rgba(34, 197, 94, 0.05) 100%);
                            border: 1px solid rgba(34, 197, 94, 0.3); border-radius: 12px; padding: 1rem;'>
                    <h4 style='color: #22c55e; margin: 0 0 0.5rem 0;'>🔥 Turmas Quase Lotadas (≥{alerta_lotado}%)</h4>
            """, unsafe_allow_html=True)
            for t in turmas_lotadas:
                st.markdown(f"""
                    <div style='background: rgba(0,0,0,0.2); border-radius: 8px; padding: 0.5rem; margin: 0.3rem 0;'>
                        <span style='color: #22c55e; font-weight: bold;'>{t['ocupacao']}%</span>
                        <span style='color: #94a3b8;'> • {t['unidade']} • {t['segmento']}</span><br>
                        <span style='color: #fff; font-size: 0.85rem;'>{t['turma'][:50]}...</span>
                    </div>
                """, unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)
        else:
            st.markdown("""
                <div style='background: rgba(251, 191, 36, 0.1); border: 1px solid rgba(251, 191, 36, 0.3);
                            border-radius: 12px; padding: 1rem; text-align: center;'>
                    <span style='color: #fbbf24; font-size: 1.2rem;'>📊 Nenhuma turma lotada ainda</span>
                </div>
            """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

# ============================================================
# HISTÓRICO DE ALERTAS
# ============================================================
with st.expander("📜 Histórico de Alertas", expanded=False):
    st.markdown("""
        <div style='background: rgba(15, 33, 55, 0.5); border-radius: 12px; padding: 1rem;'>
            <h4 style='color: #94a3b8; margin: 0 0 1rem 0;'>📊 Resumo de Alertas Atual</h4>
        </div>
    """, unsafe_allow_html=True)

    # Mostra resumo atual de alertas
    col_hist1, col_hist2, col_hist3 = st.columns(3)

    with col_hist1:
        st.markdown(f"""
            <div style='background: rgba(239, 68, 68, 0.1); border: 1px solid rgba(239, 68, 68, 0.3);
                        border-radius: 8px; padding: 1rem; text-align: center;'>
                <div style='font-size: 2.5rem; font-weight: bold; color: #ef4444;'>{total_criticas}</div>
                <div style='color: #94a3b8;'>❄️ Turmas Críticas</div>
                <div style='color: #64748b; font-size: 0.75rem;'>Ocupação &lt;{alerta_critico}%</div>
            </div>
        """, unsafe_allow_html=True)

    with col_hist2:
        st.markdown(f"""
            <div style='background: rgba(249, 115, 22, 0.1); border: 1px solid rgba(249, 115, 22, 0.3);
                        border-radius: 8px; padding: 1rem; text-align: center;'>
                <div style='font-size: 2.5rem; font-weight: bold; color: #f97316;'>{total_atencao}</div>
                <div style='color: #94a3b8;'>⚠️ Turmas Atenção</div>
                <div style='color: #64748b; font-size: 0.75rem;'>Ocupação {alerta_critico}-{alerta_atencao}%</div>
            </div>
        """, unsafe_allow_html=True)

    with col_hist3:
        st.markdown(f"""
            <div style='background: rgba(34, 197, 94, 0.1); border: 1px solid rgba(34, 197, 94, 0.3);
                        border-radius: 8px; padding: 1rem; text-align: center;'>
                <div style='font-size: 2.5rem; font-weight: bold; color: #22c55e;'>{total_lotadas}</div>
                <div style='color: #94a3b8;'>🔥 Turmas Lotadas</div>
                <div style='color: #64748b; font-size: 0.75rem;'>Ocupação ≥{alerta_lotado}%</div>
            </div>
        """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # Gráfico de histórico de alertas (se tiver histórico no banco)
    if not df_hist_total.empty and len(df_hist_total) > 1:
        st.markdown("#### 📈 Evolução da Ocupação")

        fig_hist_ocup = go.Figure()

        df_hist_total['ocupacao'] = round(df_hist_total['matriculados'] / df_hist_total['vagas'] * 100, 1)

        fig_hist_ocup.add_trace(go.Scatter(
            x=df_hist_total['data_formatada'],
            y=df_hist_total['ocupacao'],
            mode='lines+markers',
            name='Ocupação %',
            line=dict(color='#3b82f6', width=3),
            marker=dict(size=8, color='#3b82f6'),
            fill='tozeroy',
            fillcolor='rgba(59, 130, 246, 0.1)'
        ))

        # Linhas de referência para os limites de alerta
        fig_hist_ocup.add_hline(y=alerta_critico, line_dash="dash", line_color="#ef4444",
                                annotation_text=f"Crítico ({alerta_critico}%)", annotation_position="right")
        fig_hist_ocup.add_hline(y=alerta_atencao, line_dash="dash", line_color="#f97316",
                                annotation_text=f"Atenção ({alerta_atencao}%)", annotation_position="right")
        fig_hist_ocup.add_hline(y=80, line_dash="dash", line_color="#22c55e",
                                annotation_text="Meta (80%)", annotation_position="right")

        fig_hist_ocup.update_layout(
            paper_bgcolor='rgba(0,0,0,0)',
            plot_bgcolor='rgba(0,0,0,0)',
            font=dict(color='#94a3b8'),
            height=300,
            margin=dict(t=20, b=40, l=40, r=80),
            xaxis=dict(gridcolor='rgba(59, 130, 246, 0.1)', tickfont=dict(color='#94a3b8')),
            yaxis=dict(gridcolor='rgba(59, 130, 246, 0.1)', tickfont=dict(color='#94a3b8'), range=[0, 100]),
            showlegend=False
        )

        st.plotly_chart(fig_hist_ocup, use_container_width=True)
    else:
        st.info("📊 O histórico de alertas será exibido após múltiplas extrações.")

    # Lista de todas as turmas em alerta
    st.markdown("#### 📋 Lista Completa de Alertas")

    todas_turmas_alerta_sorted = sorted(todas_turmas_alerta, key=lambda x: x['ocupacao'])

    # Tabs para categorias
    tab_criticas, tab_atencao, tab_lotadas = st.tabs([f"❄️ Críticas ({total_criticas})", f"⚠️ Atenção ({total_atencao})", f"🔥 Lotadas ({total_lotadas})"])

    with tab_criticas:
        turmas_crit = [t for t in todas_turmas_alerta_sorted if t['ocupacao'] < alerta_critico]
        if turmas_crit:
            df_crit = pd.DataFrame(turmas_crit)
            df_crit = df_crit[['unidade', 'segmento', 'turma', 'vagas', 'matriculados', 'disponiveis', 'ocupacao']]
            df_crit.columns = ['Unidade', 'Segmento', 'Turma', 'Vagas', 'Matr.', 'Disp.', 'Ocup. %']
            st.dataframe(df_crit, use_container_width=True, hide_index=True)
        else:
            st.success("✅ Nenhuma turma crítica!")

    with tab_atencao:
        turmas_atenc = [t for t in todas_turmas_alerta_sorted if alerta_critico <= t['ocupacao'] < alerta_atencao]
        if turmas_atenc:
            df_atenc = pd.DataFrame(turmas_atenc)
            df_atenc = df_atenc[['unidade', 'segmento', 'turma', 'vagas', 'matriculados', 'disponiveis', 'ocupacao']]
            df_atenc.columns = ['Unidade', 'Segmento', 'Turma', 'Vagas', 'Matr.', 'Disp.', 'Ocup. %']
            st.dataframe(df_atenc, use_container_width=True, hide_index=True)
        else:
            st.success("✅ Nenhuma turma em atenção!")

    with tab_lotadas:
        turmas_lot = sorted([t for t in todas_turmas_alerta if t['ocupacao'] >= alerta_lotado], key=lambda x: -x['ocupacao'])
        if turmas_lot:
            df_lot = pd.DataFrame(turmas_lot)
            df_lot = df_lot[['unidade', 'segmento', 'turma', 'vagas', 'matriculados', 'disponiveis', 'ocupacao']]
            df_lot.columns = ['Unidade', 'Segmento', 'Turma', 'Vagas', 'Matr.', 'Disp.', 'Ocup. %']
            st.dataframe(df_lot, use_container_width=True, hide_index=True)
        else:
            st.info("📊 Nenhuma turma lotada ainda.")

st.markdown("<br>", unsafe_allow_html=True)

# ============================================================
# GAUGE DE OCUPAÇÃO GERAL
# ============================================================
st.markdown("### 🎯 Ocupação Geral")

col_gauge, col_treemap = st.columns([1, 2])

with col_gauge:
    # Gauge de ocupação
    fig_gauge = go.Figure(go.Indicator(
        mode="gauge+number+delta",
        value=ocupacao,
        number={'suffix': '%', 'font': {'size': 48, 'color': 'white'}},
        delta={'reference': 80, 'increasing': {'color': '#22c55e'}, 'decreasing': {'color': '#ef4444'}},
        gauge={
            'axis': {'range': [0, 100], 'tickwidth': 1, 'tickcolor': '#94a3b8', 'tickfont': {'color': '#94a3b8'}},
            'bar': {'color': get_ocupacao_color(ocupacao)},
            'bgcolor': 'rgba(15, 33, 55, 0.5)',
            'borderwidth': 2,
            'bordercolor': 'rgba(59, 130, 246, 0.3)',
            'steps': [
                {'range': [0, 50], 'color': 'rgba(239, 68, 68, 0.15)'},
                {'range': [50, 70], 'color': 'rgba(249, 115, 22, 0.15)'},
                {'range': [70, 80], 'color': 'rgba(251, 191, 36, 0.15)'},
                {'range': [80, 90], 'color': 'rgba(132, 204, 22, 0.15)'},
                {'range': [90, 100], 'color': 'rgba(34, 197, 94, 0.15)'}
            ],
            'threshold': {
                'line': {'color': '#ffffff', 'width': 3},
                'thickness': 0.8,
                'value': ocupacao
            }
        },
        title={'text': 'Meta: 80%', 'font': {'color': '#64748b', 'size': 14}}
    ))

    fig_gauge.update_layout(
        paper_bgcolor='rgba(0,0,0,0)',
        font={'color': '#94a3b8'},
        height=280,
        margin=dict(t=40, b=20, l=30, r=30)
    )

    st.plotly_chart(fig_gauge, use_container_width=True)

with col_treemap:
    # Treemap hierárquico
    treemap_data = []
    for unidade in vagas['unidades']:
        nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        for turma in unidade['turmas']:
            ocup_turma = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0
            treemap_data.append({
                'Unidade': nome_unidade,
                'Segmento': turma['segmento'],
                'Turma': turma['turma'][:30] + '...' if len(turma['turma']) > 30 else turma['turma'],
                'Matriculados': turma['matriculados'],
                'Vagas': turma['vagas'],
                'Ocupação': ocup_turma
            })

    df_treemap = pd.DataFrame(treemap_data)

    fig_treemap = px.treemap(
        df_treemap,
        path=['Unidade', 'Segmento', 'Turma'],
        values='Matriculados',
        color='Ocupação',
        color_continuous_scale=[
            [0, '#ef4444'],
            [0.5, '#fbbf24'],
            [0.7, '#84cc16'],
            [1, '#22c55e']
        ],
        range_color=[0, 100],
        hover_data={'Vagas': True, 'Matriculados': True, 'Ocupação': ':.1f'}
    )

    fig_treemap.update_layout(
        paper_bgcolor='rgba(0,0,0,0)',
        font={'color': '#94a3b8'},
        height=280,
        margin=dict(t=30, b=10, l=10, r=10),
        coloraxis_colorbar=dict(
            title='Ocupação %',
            tickfont=dict(color='#94a3b8'),
            titlefont=dict(color='#94a3b8')
        )
    )

    fig_treemap.update_traces(
        textfont=dict(color='white'),
        hovertemplate='<b>%{label}</b><br>Matriculados: %{value}<br>Ocupação: %{color:.1f}%<extra></extra>'
    )

    st.plotly_chart(fig_treemap, use_container_width=True)

st.markdown("<br>", unsafe_allow_html=True)

# Gráficos principais
col_left, col_right = st.columns(2)

with col_left:
    st.markdown("### 🌡️ Ocupação por Unidade")
    st.markdown("""
        <div style='display: flex; gap: 1rem; font-size: 0.75rem; margin-bottom: 0.5rem;'>
            <span style='color: #22c55e;'>🔥 90-100% Excelente</span>
            <span style='color: #84cc16;'>✨ 80-89% Muito Bom</span>
            <span style='color: #fbbf24;'>⚡ 70-79% Bom</span>
            <span style='color: #f97316;'>⚠️ 50-69% Atenção</span>
            <span style='color: #ef4444;'>❄️ &lt;50% Crítico</span>
        </div>
    """, unsafe_allow_html=True)

    df_unidades = pd.DataFrame([
        {
            'Unidade': u['nome'].split('(')[1].replace(')', '') if '(' in u['nome'] else u['nome'],
            'Ocupação': round(u['total']['matriculados'] / u['total']['vagas'] * 100, 1),
            'Matriculados': u['total']['matriculados'],
            'Vagas': u['total']['vagas']
        }
        for u in resumo['unidades']
    ])

    fig1 = go.Figure()

    # Barra de fundo (vagas totais)
    fig1.add_trace(go.Bar(
        name='Capacidade',
        x=df_unidades['Unidade'],
        y=[100] * len(df_unidades),
        marker_color='rgba(59, 130, 246, 0.12)',
        hoverinfo='skip'
    ))

    # Barra de ocupação - quanto maior, mais quente (verde)
    colors = [get_ocupacao_color(o) for o in df_unidades['Ocupação']]

    fig1.add_trace(go.Bar(
        name='Ocupação',
        x=df_unidades['Unidade'],
        y=df_unidades['Ocupação'],
        marker_color=colors,
        text=df_unidades['Ocupação'].apply(lambda x: f'{x}%'),
        textposition='outside',
        textfont=dict(color='#ffffff', size=14, family='Inter')
    ))

    fig1.update_layout(
        paper_bgcolor=PLOTLY_LAYOUT['paper_bgcolor'],
        plot_bgcolor=PLOTLY_LAYOUT['plot_bgcolor'],
        font=PLOTLY_LAYOUT['font'],
        margin=PLOTLY_LAYOUT['margin'],
        barmode='overlay',
        showlegend=False,
        height=350,
        yaxis=dict(**PLOTLY_LAYOUT['yaxis'], range=[0, 110], title=''),
        xaxis=dict(**PLOTLY_LAYOUT['xaxis'], title='')
    )

    st.plotly_chart(fig1, use_container_width=True)

with col_right:
    st.markdown("### Distribuição por Segmento")

    segmentos_total = {}
    for unidade in resumo['unidades']:
        for seg, vals in unidade['segmentos'].items():
            if seg not in segmentos_total:
                segmentos_total[seg] = {'matriculados': 0, 'vagas': 0}
            segmentos_total[seg]['matriculados'] += vals['matriculados']
            segmentos_total[seg]['vagas'] += vals['vagas']

    df_seg = pd.DataFrame([
        {'Segmento': seg, 'Matriculados': v['matriculados'], 'Vagas': v['vagas']}
        for seg, v in segmentos_total.items()
    ])

    ordem = ['Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']
    df_seg['ordem'] = df_seg['Segmento'].map({s: i for i, s in enumerate(ordem)})
    df_seg = df_seg.sort_values('ordem')

    fig2 = go.Figure()

    fig2.add_trace(go.Bar(
        name='Vagas',
        x=df_seg['Segmento'],
        y=df_seg['Vagas'],
        marker_color='rgba(59, 130, 246, 0.25)',
        text=df_seg['Vagas'],
        textposition='outside',
        textfont=dict(color='#3b82f6')
    ))

    fig2.add_trace(go.Bar(
        name='Matriculados',
        x=df_seg['Segmento'],
        y=df_seg['Matriculados'],
        marker=dict(
            color=df_seg['Matriculados'],
            colorscale=[[0, '#1e4976'], [1, '#2563eb']]
        ),
        text=df_seg['Matriculados'],
        textposition='outside',
        textfont=dict(color='#ffffff')
    ))

    fig2.update_layout(
        paper_bgcolor=PLOTLY_LAYOUT['paper_bgcolor'],
        plot_bgcolor=PLOTLY_LAYOUT['plot_bgcolor'],
        font=PLOTLY_LAYOUT['font'],
        margin=PLOTLY_LAYOUT['margin'],
        xaxis=PLOTLY_LAYOUT['xaxis'],
        yaxis=PLOTLY_LAYOUT['yaxis'],
        barmode='group',
        height=350,
        legend=dict(
            orientation='h',
            yanchor='bottom',
            y=1.02,
            xanchor='right',
            x=1,
            bgcolor='rgba(0,0,0,0)',
            font=dict(color='#94a3b8')
        )
    )

    st.plotly_chart(fig2, use_container_width=True)

st.markdown("<br>", unsafe_allow_html=True)

# ============================================================
# ANÁLISE MACRO → MICRO
# ============================================================

st.markdown("""
    <div style='background: linear-gradient(90deg, rgba(37, 99, 235, 0.2) 0%, transparent 100%);
                padding: 1rem 1.5rem; border-left: 4px solid #2563eb; border-radius: 0 12px 12px 0; margin-bottom: 1rem;'>
        <h2 style='color: #ffffff; margin: 0; font-size: 1.5rem;'>📊 Análise Detalhada</h2>
        <p style='color: #94a3b8; margin: 0.5rem 0 0 0;'>Navegue do macro ao micro: Unidades → Segmentos → Séries → Turmas</p>
    </div>
""", unsafe_allow_html=True)

# ============================================================
# NÍVEL 1: UNIDADES
# ============================================================
st.markdown("### 🏫 Nível 1: Visão por Unidade")

col_comp1, col_comp2 = st.columns(2)

with col_comp1:
    # Gráfico de Ocupação lado a lado
    df_comp_ocup = pd.DataFrame([
        {
            'Unidade': u['nome'].split('(')[1].replace(')', '') if '(' in u['nome'] else u['nome'],
            'Ocupação': round(u['total']['matriculados'] / u['total']['vagas'] * 100, 1)
        }
        for u in resumo['unidades']
    ]).sort_values('Ocupação', ascending=True)

    fig_comp1 = go.Figure()

    colors_comp = [get_ocupacao_color(o) for o in df_comp_ocup['Ocupação']]

    fig_comp1.add_trace(go.Bar(
        x=df_comp_ocup['Ocupação'],
        y=df_comp_ocup['Unidade'],
        orientation='h',
        marker_color=colors_comp,
        text=df_comp_ocup['Ocupação'].apply(lambda x: f'{x}%'),
        textposition='outside',
        textfont=dict(color='white', size=12)
    ))

    fig_comp1.update_layout(
        title=dict(text='Ocupação (%)', font=dict(color='#ffffff', size=14)),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font=dict(color='#94a3b8'),
        margin=dict(t=50, b=30, l=80, r=60),
        height=250,
        xaxis=dict(
            gridcolor='rgba(59, 130, 246, 0.1)',
            range=[0, 110],
            tickfont=dict(color='#94a3b8')
        ),
        yaxis=dict(
            tickfont=dict(color='#94a3b8')
        ),
        showlegend=False
    )

    st.plotly_chart(fig_comp1, use_container_width=True)

with col_comp2:
    # Gráfico de Matriculados x Vagas
    df_comp_matr = pd.DataFrame([
        {
            'Unidade': u['nome'].split('(')[1].replace(')', '') if '(' in u['nome'] else u['nome'],
            'Matriculados': u['total']['matriculados'],
            'Vagas': u['total']['vagas'],
            'Disponíveis': u['total']['vagas'] - u['total']['matriculados']
        }
        for u in resumo['unidades']
    ])

    fig_comp2 = go.Figure()

    fig_comp2.add_trace(go.Bar(
        name='Matriculados',
        x=df_comp_matr['Unidade'],
        y=df_comp_matr['Matriculados'],
        marker_color=COLORS['primary'],
        text=df_comp_matr['Matriculados'],
        textposition='outside',
        textfont=dict(color='white', size=11)
    ))

    fig_comp2.add_trace(go.Bar(
        name='Disponíveis',
        x=df_comp_matr['Unidade'],
        y=df_comp_matr['Disponíveis'],
        marker_color='rgba(239, 68, 68, 0.7)',
        text=df_comp_matr['Disponíveis'],
        textposition='outside',
        textfont=dict(color='#ef4444', size=11)
    ))

    fig_comp2.update_layout(
        title=dict(text='Matriculados vs Disponíveis', font=dict(color='#ffffff', size=14)),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font=dict(color='#94a3b8'),
        margin=dict(t=50, b=30, l=40, r=40),
        height=250,
        barmode='stack',
        xaxis=dict(
            tickfont=dict(color='#94a3b8')
        ),
        yaxis=dict(
            gridcolor='rgba(59, 130, 246, 0.1)',
            tickfont=dict(color='#94a3b8')
        ),
        legend=dict(
            orientation='h',
            yanchor='bottom',
            y=1.02,
            xanchor='right',
            x=1,
            bgcolor='rgba(0,0,0,0)',
            font=dict(color='#94a3b8')
        )
    )

    st.plotly_chart(fig_comp2, use_container_width=True)

st.markdown("<br>", unsafe_allow_html=True)

# 3. Tendência Histórica (Setas ↑↓)
if num_extracoes >= 2 and not df_hist_total.empty:
    st.markdown("### 📈 Tendência Recente")
    st.markdown("<p style='color: #64748b;'>Evolução desde a última extração</p>", unsafe_allow_html=True)

    # Pega as duas últimas extrações
    df_hist_sorted = df_hist_total.sort_values('data_extracao', ascending=False)

    if len(df_hist_sorted) >= 2:
        ultima = df_hist_sorted.iloc[0]
        penultima = df_hist_sorted.iloc[1]

        # Calcula deltas
        delta_matriculados = int(ultima['matriculados'] - penultima['matriculados'])
        delta_disponiveis = int(ultima['disponiveis'] - penultima['disponiveis'])
        delta_novatos = int(ultima['novatos'] - penultima['novatos'])
        delta_veteranos = int(ultima['veteranos'] - penultima['veteranos'])

        ocup_atual = round(ultima['matriculados'] / ultima['vagas'] * 100, 1)
        ocup_anterior = round(penultima['matriculados'] / penultima['vagas'] * 100, 1)
        delta_ocupacao = round(ocup_atual - ocup_anterior, 1)

        col_t1, col_t2, col_t3, col_t4, col_t5 = st.columns(5)

        def format_delta(val, inverso=False):
            """Formata delta com seta e cor"""
            if val > 0:
                return f"+{val}", "normal" if not inverso else "inverse"
            elif val < 0:
                return f"{val}", "inverse" if not inverso else "normal"
            else:
                return "0", "off"

        with col_t1:
            delta_str, delta_color = format_delta(delta_ocupacao)
            st.metric("Ocupação", f"{ocup_atual}%", delta=f"{delta_str}pp", delta_color=delta_color)

        with col_t2:
            delta_str, delta_color = format_delta(delta_matriculados)
            st.metric("Matriculados", int(ultima['matriculados']), delta=delta_str, delta_color=delta_color)

        with col_t3:
            delta_str, delta_color = format_delta(delta_disponiveis, inverso=True)
            st.metric("Disponíveis", int(ultima['disponiveis']), delta=delta_str, delta_color=delta_color)

        with col_t4:
            delta_str, delta_color = format_delta(delta_novatos)
            st.metric("Novatos", int(ultima['novatos']), delta=delta_str, delta_color=delta_color)

        with col_t5:
            delta_str, delta_color = format_delta(delta_veteranos)
            st.metric("Veteranos", int(ultima['veteranos']), delta=delta_str, delta_color=delta_color)

        # Função para gerar sparkline SVG
        def gerar_sparkline(valores, cor='#3b82f6', largura=80, altura=30):
            """Gera um SVG sparkline a partir de uma lista de valores"""
            if not valores or len(valores) < 2:
                return ''

            # Normaliza valores para caber no SVG
            min_val = min(valores)
            max_val = max(valores)
            range_val = max_val - min_val if max_val != min_val else 1

            pontos = []
            for i, v in enumerate(valores):
                x = (i / (len(valores) - 1)) * largura
                y = altura - ((v - min_val) / range_val) * (altura - 4) - 2
                pontos.append(f"{x},{y}")

            path = "M" + " L".join(pontos)

            # Determina cor baseada na tendência
            if valores[-1] > valores[0]:
                cor_linha = '#22c55e'
            elif valores[-1] < valores[0]:
                cor_linha = '#ef4444'
            else:
                cor_linha = '#94a3b8'

            svg = f'''<svg width="{largura}" height="{altura}" style="display: block; margin: 0.5rem auto;">
                <path d="{path}" fill="none" stroke="{cor_linha}" stroke-width="2" stroke-linecap="round"/>
                <circle cx="{largura}" cy="{altura - ((valores[-1] - min_val) / range_val) * (altura - 4) - 2}" r="3" fill="{cor_linha}"/>
            </svg>'''
            return svg

        # Tendência por unidade com sparklines
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("**Tendência por Unidade:**")

        df_hist_unid_sorted = df_hist_unidades.sort_values('data_extracao', ascending=True)
        unidades_unicas = df_hist_unid_sorted['unidade_nome'].unique()

        tendencias_unid = []
        for unid in unidades_unicas:
            df_u = df_hist_unid_sorted[df_hist_unid_sorted['unidade_nome'] == unid].tail(10)  # Últimas 10 extrações
            if len(df_u) >= 2:
                # Calcula ocupações históricas para sparkline
                ocupacoes = []
                for _, row in df_u.iterrows():
                    ocup = round(row['matriculados'] / row['vagas'] * 100, 1) if row['vagas'] > 0 else 0
                    ocupacoes.append(ocup)

                u_atual = df_u.iloc[-1]
                u_anterior = df_u.iloc[-2]

                ocup_u_atual = ocupacoes[-1]
                ocup_u_anterior = ocupacoes[-2]
                delta_u = round(ocup_u_atual - ocup_u_anterior, 1)

                nome_u = unid.split('(')[1].replace(')', '') if '(' in unid else unid

                if delta_u > 0:
                    seta = '↑'
                    cor = '#22c55e'
                elif delta_u < 0:
                    seta = '↓'
                    cor = '#ef4444'
                else:
                    seta = '→'
                    cor = '#94a3b8'

                tendencias_unid.append({
                    'nome': nome_u,
                    'ocupacao': ocup_u_atual,
                    'delta': delta_u,
                    'seta': seta,
                    'cor': cor,
                    'sparkline': gerar_sparkline(ocupacoes)
                })

        if tendencias_unid:
            cols_tend = st.columns(len(tendencias_unid))
            for i, t in enumerate(tendencias_unid):
                with cols_tend[i]:
                    st.markdown(f"""
                        <div style='background: linear-gradient(145deg, #0d1f35 0%, #142d4c 100%);
                                    border: 1px solid rgba(59, 130, 246, 0.2);
                                    border-radius: 12px; padding: 1rem; text-align: center;'>
                            <div style='color: #94a3b8; font-size: 0.85rem;'>{t['nome']}</div>
                            <div style='font-size: 1.5rem; font-weight: bold; color: white;'>{t['ocupacao']}%</div>
                            {t['sparkline']}
                            <div style='color: {t["cor"]}; font-size: 1rem; font-weight: bold;'>
                                {t['seta']} {'+' if t['delta'] > 0 else ''}{t['delta']}pp
                            </div>
                        </div>
                    """, unsafe_allow_html=True)

st.markdown("<br>", unsafe_allow_html=True)

# Seção de histórico
if num_extracoes >= 2:
    st.markdown("### 📈 Evolução Histórica")

    # Filtro de período
    col_periodo1, col_periodo2, col_periodo3 = st.columns([1, 1, 2])

    with col_periodo1:
        if not df_hist_total.empty:
            data_min = df_hist_total['data_extracao'].min().date()
            data_max = df_hist_total['data_extracao'].max().date()
            data_inicio = st.date_input("Data Início", value=data_min, min_value=data_min, max_value=data_max, key="hist_inicio")

    with col_periodo2:
        if not df_hist_total.empty:
            data_fim = st.date_input("Data Fim", value=data_max, min_value=data_min, max_value=data_max, key="hist_fim")

    with col_periodo3:
        st.markdown(f"<p style='color: #64748b; margin-top: 2rem;'>📅 Período: {data_inicio.strftime('%d/%m/%Y')} a {data_fim.strftime('%d/%m/%Y')}</p>", unsafe_allow_html=True)

    # Filtra dados pelo período selecionado
    df_hist_total_filtrado = df_hist_total[
        (df_hist_total['data_extracao'].dt.date >= data_inicio) &
        (df_hist_total['data_extracao'].dt.date <= data_fim)
    ]
    df_hist_unidades_filtrado = df_hist_unidades[
        (df_hist_unidades['data_extracao'].dt.date >= data_inicio) &
        (df_hist_unidades['data_extracao'].dt.date <= data_fim)
    ]

    tab1, tab2 = st.tabs(["Visão Geral", "Por Unidade"])

    with tab1:
        df_hist_total_filtrado['ocupacao'] = round(df_hist_total_filtrado['matriculados'] / df_hist_total_filtrado['vagas'] * 100, 1)

        fig_hist = go.Figure()

        fig_hist.add_trace(go.Scatter(
            x=df_hist_total_filtrado['data_formatada'],
            y=df_hist_total_filtrado['ocupacao'],
            mode='lines+markers',
            name='Ocupação',
            line=dict(color=COLORS['primary'], width=3),
            marker=dict(size=10, color=COLORS['primary']),
            fill='tozeroy',
            fillcolor='rgba(37, 99, 235, 0.1)'
        ))

        fig_hist.update_layout(
            paper_bgcolor=PLOTLY_LAYOUT['paper_bgcolor'],
            plot_bgcolor=PLOTLY_LAYOUT['plot_bgcolor'],
            font=PLOTLY_LAYOUT['font'],
            margin=PLOTLY_LAYOUT['margin'],
            xaxis=PLOTLY_LAYOUT['xaxis'],
            height=300,
            yaxis=dict(**PLOTLY_LAYOUT['yaxis'], title='Ocupação %', range=[0, 100])
        )

        st.plotly_chart(fig_hist, use_container_width=True)

    with tab2:
        fig_unid = go.Figure()
        cores_unid = [COLORS['primary'], COLORS['accent'], COLORS['info'], '#60a5fa']

        for i, unidade in enumerate(df_hist_unidades_filtrado['unidade_nome'].unique()):
            df_u = df_hist_unidades_filtrado[df_hist_unidades_filtrado['unidade_nome'] == unidade]
            nome = unidade.split('(')[1].replace(')', '') if '(' in unidade else unidade

            fig_unid.add_trace(go.Scatter(
                x=df_u['data_formatada'],
                y=df_u['matriculados'],
                mode='lines+markers',
                name=nome,
                line=dict(color=cores_unid[i % len(cores_unid)], width=2),
                marker=dict(size=8)
            ))

        fig_unid.update_layout(**PLOTLY_LAYOUT, height=300, hovermode='x unified')
        st.plotly_chart(fig_unid, use_container_width=True)

st.markdown("<br>", unsafe_allow_html=True)

# ============================================================
# ANALYTICS: RANKING, PROJEÇÃO E CRESCIMENTO
# ============================================================
st.markdown("""
    <div style='background: linear-gradient(90deg, rgba(34, 197, 94, 0.2) 0%, transparent 100%);
                padding: 1rem 1.5rem; border-left: 4px solid #22c55e; border-radius: 0 12px 12px 0; margin-bottom: 1rem;'>
        <h2 style='color: #ffffff; margin: 0; font-size: 1.5rem;'>📈 Analytics</h2>
        <p style='color: #94a3b8; margin: 0.5rem 0 0 0;'>Ranking, projeções e análise de crescimento</p>
    </div>
""", unsafe_allow_html=True)

col_analytics1, col_analytics2 = st.columns(2)

with col_analytics1:
    # Ranking das Unidades
    st.markdown("#### 🏆 Ranking de Ocupação")

    ranking_data = []
    for unidade in resumo['unidades']:
        nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        t = unidade['total']
        ocup = round(t['matriculados'] / t['vagas'] * 100, 1)
        ranking_data.append({'Unidade': nome, 'Ocupação': ocup, 'Matriculados': t['matriculados'], 'Vagas': t['vagas']})

    df_ranking = pd.DataFrame(ranking_data).sort_values('Ocupação', ascending=False)
    df_ranking['Posição'] = range(1, len(df_ranking) + 1)
    df_ranking = df_ranking[['Posição', 'Unidade', 'Ocupação', 'Matriculados', 'Vagas']]

    # Adiciona medalhas
    def add_medalha(pos):
        if pos == 1: return '🥇'
        elif pos == 2: return '🥈'
        elif pos == 3: return '🥉'
        else: return f'{pos}º'

    df_ranking['Posição'] = df_ranking['Posição'].apply(add_medalha)

    def cor_ranking(val):
        if isinstance(val, (int, float)):
            cor = get_ocupacao_color(val)
            return f'color: {cor}; font-weight: bold;'
        return ''

    styled_ranking = df_ranking.style.map(cor_ranking, subset=['Ocupação'])
    st.dataframe(styled_ranking, use_container_width=True, hide_index=True, height=200)

with col_analytics2:
    # Taxa de Crescimento
    st.markdown("#### 📊 Taxa de Crescimento")

    if num_extracoes >= 2 and not df_hist_total.empty:
        df_growth = df_hist_total.sort_values('data_extracao', ascending=True)

        if len(df_growth) >= 2:
            primeiro = df_growth.iloc[0]
            ultimo = df_growth.iloc[-1]

            crescimento_matr = ultimo['matriculados'] - primeiro['matriculados']
            taxa_crescimento = round((crescimento_matr / primeiro['matriculados']) * 100, 1) if primeiro['matriculados'] > 0 else 0

            dias = (ultimo['data_extracao'] - primeiro['data_extracao']).days
            taxa_diaria = round(crescimento_matr / dias, 1) if dias > 0 else 0

            col_g1, col_g2 = st.columns(2)
            with col_g1:
                cor_cresc = '#22c55e' if crescimento_matr >= 0 else '#ef4444'
                st.markdown(f"""
                    <div style='background: rgba(15, 33, 55, 0.8); border-radius: 12px; padding: 1rem; text-align: center;'>
                        <div style='color: #94a3b8; font-size: 0.75rem;'>CRESCIMENTO TOTAL</div>
                        <div style='color: {cor_cresc}; font-size: 1.8rem; font-weight: bold;'>
                            {'+' if crescimento_matr >= 0 else ''}{crescimento_matr}
                        </div>
                        <div style='color: #64748b; font-size: 0.7rem;'>matrículas ({taxa_crescimento}%)</div>
                    </div>
                """, unsafe_allow_html=True)
            with col_g2:
                st.markdown(f"""
                    <div style='background: rgba(15, 33, 55, 0.8); border-radius: 12px; padding: 1rem; text-align: center;'>
                        <div style='color: #94a3b8; font-size: 0.75rem;'>MÉDIA DIÁRIA</div>
                        <div style='color: #3b82f6; font-size: 1.8rem; font-weight: bold;'>
                            {'+' if taxa_diaria >= 0 else ''}{taxa_diaria}
                        </div>
                        <div style='color: #64748b; font-size: 0.7rem;'>matrículas/dia ({dias} dias)</div>
                    </div>
                """, unsafe_allow_html=True)
    else:
        st.info("Necessário mais extrações para calcular crescimento")

st.markdown("<br>", unsafe_allow_html=True)

# Projeção de Lotação
st.markdown("#### 🔮 Projeção de Lotação")

if num_extracoes >= 3 and not df_hist_total.empty:
    df_proj = df_hist_total.sort_values('data_extracao', ascending=True)

    if len(df_proj) >= 3:
        # Calcula taxa média de crescimento diário
        primeiro = df_proj.iloc[0]
        ultimo = df_proj.iloc[-1]
        dias_total = (ultimo['data_extracao'] - primeiro['data_extracao']).days

        if dias_total > 0:
            taxa_diaria_matr = (ultimo['matriculados'] - primeiro['matriculados']) / dias_total

            projecoes = []
            for unidade in resumo['unidades']:
                nome = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
                t = unidade['total']
                vagas_disp = t['vagas'] - t['matriculados']

                if taxa_diaria_matr > 0 and vagas_disp > 0:
                    # Proporção da unidade no total
                    prop = t['matriculados'] / total['matriculados'] if total['matriculados'] > 0 else 0.25
                    taxa_unidade = taxa_diaria_matr * prop

                    if taxa_unidade > 0:
                        dias_lotacao = int(vagas_disp / taxa_unidade)
                        if dias_lotacao < 365:  # Só mostra se for menos de 1 ano
                            from datetime import timedelta
                            data_lotacao = datetime.now() + timedelta(days=dias_lotacao)
                            projecoes.append({
                                'Unidade': nome,
                                'Disponíveis': vagas_disp,
                                'Taxa/dia': round(taxa_unidade, 2),
                                'Dias p/ Lotar': dias_lotacao,
                                'Previsão': data_lotacao.strftime('%d/%m/%Y')
                            })

            if projecoes:
                df_proj_display = pd.DataFrame(projecoes).sort_values('Dias p/ Lotar')

                def cor_dias(val):
                    if val <= 30: return 'color: #22c55e; font-weight: bold;'
                    elif val <= 90: return 'color: #84cc16;'
                    elif val <= 180: return 'color: #fbbf24;'
                    else: return 'color: #94a3b8;'

                styled_proj = df_proj_display.style.map(cor_dias, subset=['Dias p/ Lotar'])
                st.dataframe(styled_proj, use_container_width=True, hide_index=True)
            else:
                st.info("Nenhuma unidade com projeção de lotação em até 1 ano")
        else:
            st.info("Período muito curto para projeção")
else:
    st.info("Necessário pelo menos 3 extrações para calcular projeções")

st.markdown("<br>", unsafe_allow_html=True)

# ============================================================
# NÍVEL 2: SEGMENTOS
# ============================================================
st.markdown("### 📚 Nível 2: Visão por Segmento")

# Visão geral dos segmentos (todas unidades)
segmentos_geral = {}
for unidade in resumo['unidades']:
    for seg, vals in unidade['segmentos'].items():
        if seg not in segmentos_geral:
            segmentos_geral[seg] = {'vagas': 0, 'matriculados': 0, 'novatos': 0, 'veteranos': 0}
        segmentos_geral[seg]['vagas'] += vals['vagas']
        segmentos_geral[seg]['matriculados'] += vals['matriculados']
        segmentos_geral[seg]['novatos'] += vals['novatos']
        segmentos_geral[seg]['veteranos'] += vals['veteranos']

# Cards dos segmentos
col_seg1, col_seg2, col_seg3, col_seg4 = st.columns(4)
segmentos_ordem_cards = ['Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']

for i, (col, seg) in enumerate(zip([col_seg1, col_seg2, col_seg3, col_seg4], segmentos_ordem_cards)):
    if seg in segmentos_geral:
        dados = segmentos_geral[seg]
        ocup_seg = round(dados['matriculados'] / dados['vagas'] * 100, 1) if dados['vagas'] > 0 else 0
        cor = get_ocupacao_color(ocup_seg)

        with col:
            st.markdown(f"""
                <div style='background: linear-gradient(145deg, #0d1f35 0%, #142d4c 100%);
                            border: 1px solid rgba(59, 130, 246, 0.2);
                            border-radius: 12px; padding: 1rem; text-align: center;'>
                    <div style='color: #94a3b8; font-size: 0.8rem; text-transform: uppercase;'>{seg}</div>
                    <div style='font-size: 2rem; font-weight: bold; color: {cor};'>{ocup_seg}%</div>
                    <div style='color: #64748b; font-size: 0.75rem;'>{dados['matriculados']}/{dados['vagas']} vagas</div>
                </div>
            """, unsafe_allow_html=True)

st.markdown("<br>", unsafe_allow_html=True)

# Filtro para detalhar segmento
col_seg_filter, col_search = st.columns([1, 2])

with col_seg_filter:
    segmentos_disponiveis = ['Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']
    segmento_filtro = st.selectbox("🔍 Detalhar Segmento", segmentos_disponiveis, key="filtro_segmento")

with col_search:
    busca_turma = st.text_input("🔍 Buscar Turma/Série", placeholder="Digite o nome...", key="busca_turma")

# Dados do segmento selecionado em todas as unidades
dados_segmento_todas = []
for unidade in resumo['unidades']:
    nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
    if segmento_filtro in unidade['segmentos']:
        seg_data = unidade['segmentos'][segmento_filtro]
        disponíveis = seg_data['vagas'] - seg_data['matriculados']
        ocup = round(seg_data['matriculados'] / seg_data['vagas'] * 100, 1) if seg_data['vagas'] > 0 else 0

        if ocup >= 90: status = '🔥 Excelente'
        elif ocup >= 80: status = '✨ Muito Bom'
        elif ocup >= 70: status = '⚡ Bom'
        elif ocup >= 50: status = '⚠️ Atenção'
        else: status = '❄️ Crítico'

        # Calcula pré-matriculados
        unidade_vagas_data = next((u for u in vagas['unidades'] if u['codigo'] == unidade['codigo']), None)
        pre_matr = sum(t['pre_matriculados'] for t in unidade_vagas_data['turmas'] if t['segmento'] == segmento_filtro) if unidade_vagas_data else 0

        dados_segmento_todas.append({
            'Unidade': nome_unidade,
            'Vagas': seg_data['vagas'],
            'Novatos': seg_data['novatos'],
            'Veteranos': seg_data['veteranos'],
            'Matriculados': seg_data['matriculados'],
            'Disponíveis': disponíveis,
            'Ocupação %': ocup,
            'Status': status,
            'Pré-Matr.': pre_matr
        })

df_seg_todas = pd.DataFrame(dados_segmento_todas)

# Estilização
def barra_ocup_todas(val):
    if val >= 90: cor = '#22c55e'
    elif val >= 80: cor = '#84cc16'
    elif val >= 70: cor = '#fbbf24'
    elif val >= 50: cor = '#f97316'
    else: cor = '#ef4444'
    return f'background: linear-gradient(90deg, {cor} {val}%, transparent {val}%); color: white; font-weight: bold;'

def colorir_status_todas(val):
    base = 'font-weight: 600; font-family: "SF Pro Display", system-ui, sans-serif; letter-spacing: 0.5px; text-transform: uppercase; font-size: 11px;'
    if 'Excelente' in val: return f'{base} color: #22c55e;'
    elif 'Muito Bom' in val: return f'{base} color: #84cc16;'
    elif 'Bom' in val: return f'{base} color: #fbbf24;'
    elif 'Atenção' in val: return f'{base} color: #f97316;'
    else: return f'{base} color: #ef4444;'

st.markdown(f"**{segmento_filtro}** em todas as unidades:")
styled_seg_todas = df_seg_todas.style.map(barra_ocup_todas, subset=['Ocupação %']).map(colorir_status_todas, subset=['Status'])
st.dataframe(styled_seg_todas, use_container_width=True, hide_index=True)

# Busca de turmas
if busca_turma:
    st.markdown(f"### 🔍 Resultados para: *{busca_turma}*")
    turmas_encontradas = []
    for unidade in vagas['unidades']:
        nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
        for turma in unidade['turmas']:
            if busca_turma.lower() in turma['turma'].lower():
                disponíveis = turma['vagas'] - turma['matriculados']
                ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0

                if ocup >= 90: status = '🔥 Excelente'
                elif ocup >= 80: status = '✨ Muito Bom'
                elif ocup >= 70: status = '⚡ Bom'
                elif ocup >= 50: status = '⚠️ Atenção'
                else: status = '❄️ Crítico'

                turmas_encontradas.append({
                    'Unidade': nome_unidade,
                    'Segmento': turma['segmento'],
                    'Turma': turma['turma'],
                    'Vagas': turma['vagas'],
                    'Matriculados': turma['matriculados'],
                    'Disponíveis': disponíveis,
                    'Ocupação %': ocup,
                    'Status': status
                })

    if turmas_encontradas:
        df_busca = pd.DataFrame(turmas_encontradas)
        styled_busca = df_busca.style.map(barra_ocup_todas, subset=['Ocupação %']).map(colorir_status_todas, subset=['Status'])
        st.dataframe(styled_busca, use_container_width=True, hide_index=True)
    else:
        st.info(f"Nenhuma turma encontrada com '{busca_turma}'")

st.markdown("<br>", unsafe_allow_html=True)

# ============================================================
# NÍVEL 3: SÉRIES
# ============================================================
st.markdown("### 📖 Nível 3: Visão por Série")

# Função para extrair série do nome da turma
def extrair_serie(nome_turma):
    """Extrai o nome da série do nome completo da turma"""
    if ' - Turma' in nome_turma:
        return nome_turma.split(' - Turma')[0].strip()
    elif ' - ' in nome_turma:
        partes = nome_turma.split(' - ')
        return partes[0].strip() if len(partes) > 0 else nome_turma
    return nome_turma

# Normaliza nome da série para agrupar variações
def normalizar_serie(serie):
    """Normaliza variações de nomes de séries"""
    serie_lower = serie.lower()

    # Ed. Infantil
    if 'infantil ii' in serie_lower: return 'Infantil II'
    if 'infantil iii' in serie_lower: return 'Infantil III'
    if 'infantil iv' in serie_lower: return 'Infantil IV'
    if 'infantil v' in serie_lower: return 'Infantil V'

    # Fund. I
    if '1º ano' in serie_lower or '1 ano' in serie_lower: return '1º Ano'
    if '2º ano' in serie_lower or '2 ano' in serie_lower or '2ºano' in serie_lower: return '2º Ano'
    if '3º ano' in serie_lower or '3 ano' in serie_lower: return '3º Ano'
    if '4º ano' in serie_lower or '4 ano' in serie_lower or '4 ºano' in serie_lower: return '4º Ano'
    if '5º ano' in serie_lower or '5 ano' in serie_lower: return '5º Ano'

    # Fund. II
    if '6º ano' in serie_lower or '6 ano' in serie_lower: return '6º Ano'
    if '7º ano' in serie_lower or '7 ano' in serie_lower: return '7º Ano'
    if '8º ano' in serie_lower or '8 ano' in serie_lower: return '8º Ano'
    if '9º ano' in serie_lower or '9 ano' in serie_lower: return '9º Ano'

    # Ens. Médio
    if '1ª série' in serie_lower or '1º ano médio' in serie_lower: return '1ª Série EM'
    if '2ª série' in serie_lower or '2º ano médio' in serie_lower: return '2ª Série EM'
    if '3ª série' in serie_lower or '3º ano' in serie_lower and 'médio' in serie_lower: return '3ª Série EM'

    return serie

# Agrupa turmas por série no segmento selecionado
series_data = {}
for unidade in vagas['unidades']:
    nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']
    for turma in unidade['turmas']:
        if turma['segmento'] == segmento_filtro:
            serie_original = extrair_serie(turma['turma'])
            serie = normalizar_serie(serie_original)

            if serie not in series_data:
                series_data[serie] = {'vagas': 0, 'matriculados': 0, 'novatos': 0, 'veteranos': 0, 'pre_matriculados': 0, 'turmas': 0}

            series_data[serie]['vagas'] += turma['vagas']
            series_data[serie]['matriculados'] += turma['matriculados']
            series_data[serie]['novatos'] += turma['novatos']
            series_data[serie]['veteranos'] += turma['veteranos']
            series_data[serie]['pre_matriculados'] += turma['pre_matriculados']
            series_data[serie]['turmas'] += 1

# Monta DataFrame das séries
if series_data:
    dados_series = []
    for serie, vals in series_data.items():
        ocup = round(vals['matriculados'] / vals['vagas'] * 100, 1) if vals['vagas'] > 0 else 0

        if ocup >= 90: status = '🔥 Excelente'
        elif ocup >= 80: status = '✨ Muito Bom'
        elif ocup >= 70: status = '⚡ Bom'
        elif ocup >= 50: status = '⚠️ Atenção'
        else: status = '❄️ Crítico'

        dados_series.append({
            'Série': serie,
            'Turmas': vals['turmas'],
            'Vagas': vals['vagas'],
            'Novatos': vals['novatos'],
            'Veteranos': vals['veteranos'],
            'Matriculados': vals['matriculados'],
            'Disponíveis': vals['vagas'] - vals['matriculados'],
            'Ocupação %': ocup,
            'Status': status,
            'Pré-Matr.': vals['pre_matriculados']
        })

    df_series = pd.DataFrame(dados_series)
    df_series = df_series.sort_values('Ocupação %', ascending=True)

    st.markdown(f"**Séries do segmento {segmento_filtro}:**")

    # Estilização
    def barra_ocup_series(val):
        if val >= 90: cor = '#22c55e'
        elif val >= 80: cor = '#84cc16'
        elif val >= 70: cor = '#fbbf24'
        elif val >= 50: cor = '#f97316'
        else: cor = '#ef4444'
        return f'background: linear-gradient(90deg, {cor} {val}%, transparent {val}%); color: white; font-weight: bold;'

    def colorir_status_series(val):
        base = 'font-weight: 600; font-family: "SF Pro Display", system-ui, sans-serif; letter-spacing: 0.5px; text-transform: uppercase; font-size: 11px;'
        if 'Excelente' in val: return f'{base} color: #22c55e;'
        elif 'Muito Bom' in val: return f'{base} color: #84cc16;'
        elif 'Bom' in val: return f'{base} color: #fbbf24;'
        elif 'Atenção' in val: return f'{base} color: #f97316;'
        else: return f'{base} color: #ef4444;'

    styled_series = df_series.style.map(barra_ocup_series, subset=['Ocupação %']).map(colorir_status_series, subset=['Status'])
    st.dataframe(styled_series, use_container_width=True, hide_index=True)

    # Gráfico de barras das séries
    fig_series = go.Figure()

    colors_series = [get_ocupacao_color(o) for o in df_series['Ocupação %']]

    fig_series.add_trace(go.Bar(
        x=df_series['Série'],
        y=df_series['Ocupação %'],
        marker_color=colors_series,
        text=df_series['Ocupação %'].apply(lambda x: f'{x}%'),
        textposition='outside',
        textfont=dict(color='white', size=11)
    ))

    fig_series.update_layout(
        title=dict(text=f'Ocupação por Série - {segmento_filtro}', font=dict(color='#ffffff', size=14)),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font=dict(color='#94a3b8'),
        margin=dict(t=50, b=40, l=40, r=40),
        height=300,
        xaxis=dict(tickfont=dict(color='#94a3b8')),
        yaxis=dict(gridcolor='rgba(59, 130, 246, 0.1)', tickfont=dict(color='#94a3b8'), range=[0, 110])
    )

    st.plotly_chart(fig_series, use_container_width=True)
else:
    st.info(f"Nenhuma série encontrada para {segmento_filtro}")

st.markdown("<br>", unsafe_allow_html=True)

# ============================================================
# NÍVEL 4: TURMAS
# ============================================================
st.markdown("### 🎓 Nível 4: Visão por Turma")

# Filtros para turmas
col_turma_unid, col_turma_seg = st.columns(2)

with col_turma_unid:
    unidades_nomes_turma = ['Todas'] + [u['nome'].split('(')[1].replace(')', '') if '(' in u['nome'] else u['nome'] for u in resumo['unidades']]
    unidade_turma = st.selectbox("Filtrar por Unidade", unidades_nomes_turma, key="turma_unidade")

with col_turma_seg:
    segmentos_turma = ['Todos', 'Ed. Infantil', 'Fund. I', 'Fund. II', 'Ens. Médio']
    segmento_turma = st.selectbox("Filtrar por Segmento", segmentos_turma, index=segmentos_turma.index(segmento_filtro) if segmento_filtro in segmentos_turma else 0, key="turma_segmento")

# Coleta todas as turmas
todas_turmas_nivel4 = []
for unidade in vagas['unidades']:
    nome_unidade = unidade['nome'].split('(')[1].replace(')', '') if '(' in unidade['nome'] else unidade['nome']

    # Filtra por unidade
    if unidade_turma != 'Todas' and nome_unidade != unidade_turma:
        continue

    for turma in unidade['turmas']:
        # Filtra por segmento
        if segmento_turma != 'Todos' and turma['segmento'] != segmento_turma:
            continue

        serie = normalizar_serie(extrair_serie(turma['turma']))
        disponíveis = turma['vagas'] - turma['matriculados']
        ocup = round(turma['matriculados'] / turma['vagas'] * 100, 1) if turma['vagas'] > 0 else 0

        if ocup >= 90: status = '🔥 Excelente'
        elif ocup >= 80: status = '✨ Muito Bom'
        elif ocup >= 70: status = '⚡ Bom'
        elif ocup >= 50: status = '⚠️ Atenção'
        else: status = '❄️ Crítico'

        todas_turmas_nivel4.append({
            'Unidade': nome_unidade,
            'Segmento': turma['segmento'],
            'Série': serie,
            'Turma': turma['turma'],
            'Vagas': turma['vagas'],
            'Novatos': turma['novatos'],
            'Veteranos': turma['veteranos'],
            'Matriculados': turma['matriculados'],
            'Disponíveis': disponíveis,
            'Ocupação %': ocup,
            'Status': status,
            'Pré-Matr.': turma['pre_matriculados']
        })

if todas_turmas_nivel4:
    df_turmas_nivel4 = pd.DataFrame(todas_turmas_nivel4)
    df_turmas_nivel4 = df_turmas_nivel4.sort_values('Ocupação %', ascending=True)

    st.markdown(f"**{len(df_turmas_nivel4)} turmas encontradas:**")

    # Estilização
    def barra_ocup_turma(val):
        if val >= 90: cor = '#22c55e'
        elif val >= 80: cor = '#84cc16'
        elif val >= 70: cor = '#fbbf24'
        elif val >= 50: cor = '#f97316'
        else: cor = '#ef4444'
        return f'background: linear-gradient(90deg, {cor} {val}%, transparent {val}%); color: white; font-weight: bold;'

    def colorir_status_turma(val):
        base = 'font-weight: 600; font-family: "SF Pro Display", system-ui, sans-serif; letter-spacing: 0.5px; text-transform: uppercase; font-size: 11px;'
        if 'Excelente' in val: return f'{base} color: #22c55e;'
        elif 'Muito Bom' in val: return f'{base} color: #84cc16;'
        elif 'Bom' in val: return f'{base} color: #fbbf24;'
        elif 'Atenção' in val: return f'{base} color: #f97316;'
        else: return f'{base} color: #ef4444;'

    styled_turmas = df_turmas_nivel4.style.map(barra_ocup_turma, subset=['Ocupação %']).map(colorir_status_turma, subset=['Status'])
    st.dataframe(styled_turmas, use_container_width=True, hide_index=True, height=400)
else:
    st.info("Nenhuma turma encontrada com os filtros selecionados")

st.markdown("<br>", unsafe_allow_html=True)

# Legenda de Status destacada
st.markdown("""
    <div style='display: flex; justify-content: center; gap: 1.5rem; padding: 1rem; background: rgba(15, 33, 55, 0.5); border-radius: 12px; margin: 1rem 0;'>
        <span style='color: #22c55e; font-weight: 600;'>🔥 EXCELENTE (90-100%)</span>
        <span style='color: #84cc16; font-weight: 600;'>✨ MUITO BOM (80-89%)</span>
        <span style='color: #fbbf24; font-weight: 600;'>⚡ BOM (70-79%)</span>
        <span style='color: #f97316; font-weight: 600;'>⚠️ ATENÇÃO (50-69%)</span>
        <span style='color: #ef4444; font-weight: 600;'>❄️ CRÍTICO (&lt;50%)</span>
    </div>
""", unsafe_allow_html=True)

# Footer
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown(f"""
    <div style='text-align: center; color: #64748b; font-size: 0.8rem; padding: 2rem 0;'>
        <p>Dashboard atualizado automaticamente às 6h • Última extração: {resumo['data_extracao'][:16].replace('T', ' ')}</p>
        <p style='color: #475569;'>Colégio Elo © 2026</p>
    </div>
""", unsafe_allow_html=True)
